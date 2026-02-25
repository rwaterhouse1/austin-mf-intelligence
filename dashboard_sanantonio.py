"""
San Antonio Multifamily Intelligence Dashboard
===============================================
Run with: streamlit run dashboard_sanantonio.py

Data source: data.sanantonio.gov (CKAN) via pipeline_sanantonio.py
Units are estimated from AREA (SF) / 900 — no unit count field in SA data.

Tabs:
  1. Market Overview    — Supply pressure score by submarket
  2. Supply Pipeline    — Under construction + projected delivery dates
  3. Absorption         — Delivery volume vs vacancy over time
  4. Timing Intelligence — Buy/sell signals by submarket
  5. Permit Browser     — Raw permit data, searchable
"""

import io
import os
from datetime import datetime, timedelta
from typing import Optional

import pandas as pd
import plotly.graph_objects as go
import streamlit as st
import psycopg2
import psycopg2.extras
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

load_dotenv()

# Streamlit Cloud secrets → env var → local default
try:
    DB_DSN = st.secrets["SANANTONIO_DATABASE_URL"]
except (KeyError, FileNotFoundError, AttributeError):
    DB_DSN = os.getenv("SANANTONIO_DATABASE_URL", "postgresql://postgres:postgres@localhost:5432/sanantonio_permits")

# ─────────────────────────────────────────────────────────────────────────────
# COSTAR DATA — San Antonio MSA submarkets
# Update these from your GeographyList.xlsx periodically or use the uploader.
# Submarket names must match the ZIP_CROSSWALK in pipeline_sanantonio.py.
# Estimates based on SA market conditions as of early 2025.
# ─────────────────────────────────────────────────────────────────────────────
COSTAR_DATA = {
    "Downtown": {
        "vacancy": 0.148,
        "rent_growth": 0.008,
        "inventory": 8200,
        "under_constr": 620,
        "delivered_12mo": 380,
        "asking_rent": 1850,
        "absorption_12mo": 310,
        "avg_days_on_market": 58,
        "concession_pct": 0.055,
    },
    "North Central": {
        "vacancy": 0.112,
        "rent_growth": -0.018,
        "inventory": 14500,
        "under_constr": 290,
        "delivered_12mo": 210,
        "asking_rent": 1420,
        "absorption_12mo": 390,
        "avg_days_on_market": 46,
        "concession_pct": 0.035,
    },
    "Northwest": {
        "vacancy": 0.106,
        "rent_growth": -0.022,
        "inventory": 22800,
        "under_constr": 780,
        "delivered_12mo": 540,
        "asking_rent": 1310,
        "absorption_12mo": 620,
        "avg_days_on_market": 44,
        "concession_pct": 0.040,
    },
    "Northeast": {
        "vacancy": 0.131,
        "rent_growth": -0.035,
        "inventory": 18600,
        "under_constr": 410,
        "delivered_12mo": 460,
        "asking_rent": 1175,
        "absorption_12mo": 520,
        "avg_days_on_market": 52,
        "concession_pct": 0.048,
    },
    "South": {
        "vacancy": 0.119,
        "rent_growth": -0.024,
        "inventory": 11200,
        "under_constr": 180,
        "delivered_12mo": 120,
        "asking_rent": 1090,
        "absorption_12mo": 195,
        "avg_days_on_market": 50,
        "concession_pct": 0.038,
    },
    "Southeast": {
        "vacancy": 0.142,
        "rent_growth": -0.041,
        "inventory": 9400,
        "under_constr": 95,
        "delivered_12mo": 80,
        "asking_rent": 1045,
        "absorption_12mo": 115,
        "avg_days_on_market": 57,
        "concession_pct": 0.045,
    },
    "Southwest": {
        "vacancy": 0.123,
        "rent_growth": -0.028,
        "inventory": 16800,
        "under_constr": 560,
        "delivered_12mo": 390,
        "asking_rent": 1155,
        "absorption_12mo": 480,
        "avg_days_on_market": 48,
        "concession_pct": 0.042,
    },
    "West": {
        "vacancy": 0.135,
        "rent_growth": -0.032,
        "inventory": 7600,
        "under_constr": 120,
        "delivered_12mo": 95,
        "asking_rent": 1010,
        "absorption_12mo": 110,
        "avg_days_on_market": 55,
        "concession_pct": 0.050,
    },
    "Medical Center": {
        "vacancy": 0.095,
        "rent_growth": 0.012,
        "inventory": 12400,
        "under_constr": 340,
        "delivered_12mo": 180,
        "asking_rent": 1480,
        "absorption_12mo": 350,
        "avg_days_on_market": 38,
        "concession_pct": 0.028,
    },
    "Stone Oak": {
        "vacancy": 0.088,
        "rent_growth": 0.021,
        "inventory": 19200,
        "under_constr": 920,
        "delivered_12mo": 610,
        "asking_rent": 1620,
        "absorption_12mo": 840,
        "avg_days_on_market": 34,
        "concession_pct": 0.020,
    },
    "Helotes/Leon Valley": {
        "vacancy": 0.097,
        "rent_growth": -0.008,
        "inventory": 7800,
        "under_constr": 310,
        "delivered_12mo": 240,
        "asking_rent": 1265,
        "absorption_12mo": 290,
        "avg_days_on_market": 40,
        "concession_pct": 0.030,
    },
    "Schertz/Cibolo": {
        "vacancy": 0.114,
        "rent_growth": -0.045,
        "inventory": 8900,
        "under_constr": 480,
        "delivered_12mo": 520,
        "asking_rent": 1340,
        "absorption_12mo": 580,
        "avg_days_on_market": 49,
        "concession_pct": 0.052,
    },
    "New Braunfels": {
        "vacancy": 0.162,
        "rent_growth": -0.058,
        "inventory": 10100,
        "under_constr": 690,
        "delivered_12mo": 750,
        "asking_rent": 1295,
        "absorption_12mo": 810,
        "avg_days_on_market": 66,
        "concession_pct": 0.068,
    },
}

# ─────────────────────────────────────────────────────────────────────────────
# MATTHEWS BRAND COLORS
# ─────────────────────────────────────────────────────────────────────────────
ACCENT    = "#C8102E"   # Matthews red — highlights, active tabs, key metrics
NAVY      = "#1A1A2E"   # Deep navy/charcoal — headers and primary text
GREEN     = "#16A34A"   # Success / BUY
AMBER     = "#D97706"   # Warning / HOLD
RED       = "#C8102E"   # Danger / SELL (same as accent)
BG        = "#FFFFFF"   # Primary background — white
CARD_BG   = "#F8F9FA"   # Card background — light gray
BORDER    = "#E5E7EB"   # Subtle gray border
TEXT      = "#1A1A2E"   # Primary text — navy
MUTED     = "#6B7280"   # Secondary text — gray

PLOTLY_LAYOUT = dict(
    paper_bgcolor=BG,
    plot_bgcolor=BG,
    font=dict(color=TEXT, family="'DM Sans', sans-serif"),
    xaxis=dict(gridcolor=BORDER, zerolinecolor=BORDER, tickfont=dict(color=MUTED)),
    yaxis=dict(gridcolor=BORDER, zerolinecolor=BORDER, tickfont=dict(color=MUTED)),
    legend=dict(bgcolor="rgba(255,255,255,0.9)", bordercolor=BORDER, font=dict(color=TEXT)),
    margin=dict(t=40, r=20, b=40, l=60),
)

# ─────────────────────────────────────────────────────────────────────────────
# PAGE CONFIG
# ─────────────────────────────────────────────────────────────────────────────
st.set_page_config(page_title="San Antonio MF Intelligence", page_icon="🏢", layout="wide")

st.markdown(f"""
<style>
@import url('https://fonts.googleapis.com/css2?family=DM+Mono:wght@300;400;500&family=Inter:wght@300;400;500;600;700&family=DM+Sans:wght@300;400;500;600&display=swap');

html, body, [class*="css"] {{ background-color: {BG}; color: {TEXT}; font-family: 'DM Sans', sans-serif; }}
.main {{ background-color: {BG}; }}
.stApp {{ margin-top: 0; padding-top: 0; }}
.block-container {{ padding-top: 0 !important; padding-right: 1rem; padding-bottom: 1.5rem; padding-left: 1rem; max-width: 1600px; }}
header[data-testid="stHeader"] {{ display: none !important; }}
.stDeployButton {{ display: none !important; }}
#MainMenu {{ display: none !important; }}
div[data-testid="stToolbar"] {{ display: none !important; }}
div[data-testid="stDecoration"] {{ display: none !important; }}
.uploadedFile {{ display: none !important; }}
section[data-testid="stFileUploadDropzone"] {{ display: none !important; }}
div[data-testid="stFileDropzoneInput"] {{ display: none !important; }}
[data-testid="stFileUploader"] {{ display: none !important; }}
.stFileDropzone {{ display: none !important; }}
div[data-testid="stAppViewBlockContainer"] [data-testid="stFileUploader"] {{ display: none !important; }}
.drag-drop-container {{ display: none !important; }}
[data-testid="stFileDropzone"] {{ display: none !important; }}
.dash-header {{ font-family: 'Inter', sans-serif; font-size: clamp(1.2rem, 3vw, 2.2rem); font-weight: 700; letter-spacing: 0.02em; color: {NAVY}; line-height: 1.1; white-space: nowrap; overflow: hidden; text-overflow: ellipsis; }}
.dash-header span {{ color: {ACCENT}; }}
.dash-sub {{ font-family: 'DM Mono', monospace; font-size: 0.72rem; color: {MUTED}; letter-spacing: 0.18em; text-transform: uppercase; margin-bottom: 1.5rem; }}
.kpi-card {{ background: {CARD_BG}; border: 1px solid {BORDER}; border-left: 3px solid {ACCENT}; border-radius: 4px; padding: 1rem 1.2rem; margin-bottom: 1rem; }}
.kpi-label {{ font-family: 'DM Mono', monospace; font-size: 0.65rem; color: {MUTED}; letter-spacing: 0.15em; text-transform: uppercase; margin-bottom: 0.4rem; }}
.kpi-value {{ font-family: 'Inter', sans-serif; font-size: 1.9rem; font-weight: 700; color: {NAVY}; line-height: 1; }}
.section-title {{ font-family: 'DM Mono', monospace; font-size: 0.65rem; color: {MUTED}; letter-spacing: 0.2em; text-transform: uppercase; border-bottom: 1px solid {BORDER}; padding-bottom: 0.5rem; margin-bottom: 1rem; }}
.stTabs [data-baseweb="tab-list"] {{ background-color: {BG}; border-bottom: 1px solid {BORDER}; gap: 0; }}
.stTabs [data-baseweb="tab"] {{ font-family: 'DM Mono', monospace; font-size: 0.7rem; letter-spacing: 0.12em; text-transform: uppercase; color: {MUTED}; padding: 0.75rem 1.5rem; border-bottom: 2px solid transparent; background: transparent; }}
.stTabs [aria-selected="true"] {{ color: {ACCENT} !important; border-bottom: 2px solid {ACCENT} !important; background: transparent !important; }}
.stDataFrame {{ border: 1px solid {BORDER}; border-radius: 4px; }}
input, select, textarea {{ background-color: {CARD_BG} !important; color: {TEXT} !important; border: 1px solid {BORDER} !important; border-radius: 4px !important; }}
::-webkit-scrollbar {{ width: 4px; height: 4px; }}
::-webkit-scrollbar-track {{ background: {BG}; }}
::-webkit-scrollbar-thumb {{ background: {BORDER}; border-radius: 2px; }}
</style>
""", unsafe_allow_html=True)

# ─────────────────────────────────────────────────────────────────────────────
# DATA LOADING
# ─────────────────────────────────────────────────────────────────────────────
@st.cache_data(ttl=300)
def load_permits(min_year=2010):
    conn = psycopg2.connect(DB_DSN)
    df = pd.read_sql(f"""
        SELECT permit_num, issue_date, submitted_date, address, zip_code,
               latitude, longitude, area_sf, total_units, project_name,
               work_class, cd, submarket_name,
               delivery_year, delivery_quarter, delivery_yyyyq
        FROM sa_projects
        WHERE total_units >= 5 AND delivery_year >= {min_year} AND issue_date IS NOT NULL
        ORDER BY issue_date DESC
    """, conn)
    conn.close()
    df["issue_date"] = pd.to_datetime(df["issue_date"])
    return df

@st.cache_data(ttl=300)
def load_quarterly():
    conn = psycopg2.connect(DB_DSN)
    df = pd.read_sql("""
        SELECT submarket_name, delivery_year, delivery_quarter,
               delivery_yyyyq, project_count, total_units_delivered
        FROM sa_submarket_deliveries WHERE delivery_year >= 2010 ORDER BY delivery_yyyyq
    """, conn)
    conn.close()
    return df

def get_costar_df():
    rows = [{"submarket_name": k, **v} for k, v in COSTAR_DATA.items()]
    return pd.DataFrame(rows)

# ─────────────────────────────────────────────────────────────────────────────
# PRESSURE SCORE + SIGNALS
# ─────────────────────────────────────────────────────────────────────────────
def pressure_score(row):
    # Vacancy: 25 pts — higher vacancy = more pressure
    v = min((row["vacancy"] - 0.08) / 0.15, 1.0) * 25

    # Deliveries vs inventory: 20 pts
    d = min(row["delivered_12mo"] / max(row["inventory"], 1) / 0.12, 1.0) * 20

    # Under construction vs inventory: 20 pts
    u = min(row["under_constr"] / max(row["inventory"], 1) / 0.15, 1.0) * 20

    # Rent growth (inverted): 15 pts — negative growth = pressure
    r = min((-row["rent_growth"]) / 0.08, 1.0) * 15

    # Absorption vs deliveries: 10 pts — low absorption = pressure
    absorption_ratio = row.get("absorption_12mo", 0) / max(row["delivered_12mo"], 1)
    a = max(0.0, min((1 - absorption_ratio) / 0.5, 1.0)) * 10

    # Days on market: 5 pts — above 45 days = pressure
    dom = row.get("avg_days_on_market", 45)
    dom_score = max(0.0, min((dom - 45) / 60, 1.0)) * 5

    # Concessions: 5 pts — above 4% = distress signal
    conc = row.get("concession_pct", 0)
    conc_score = max(0.0, min(max(conc - 0.04, 0) / 0.10, 1.0)) * 5

    return round(max(0, v + d + u + r + a + dom_score + conc_score), 1)

def sig(score):
    if score >= 60: return "SELL"
    if score >= 35: return "HOLD"
    return "BUY"

def sig_color(score):
    if score >= 60: return RED
    if score >= 35: return AMBER
    return GREEN

# ─────────────────────────────────────────────────────────────────────────────
# LOAD DATA
# ─────────────────────────────────────────────────────────────────────────────
try:
    df = load_permits()
    dq = load_quarterly()
    db_ok = True
except Exception as e:
    st.error(f"DB error: {e}")
    df = pd.DataFrame()
    dq = pd.DataFrame()
    db_ok = False

dc = get_costar_df()
dc["score"] = dc.apply(pressure_score, axis=1)
dc["signal"] = dc["score"].apply(sig)

# ─────────────────────────────────────────────────────────────────────────────
# HEADER
# ─────────────────────────────────────────────────────────────────────────────
h1, h2 = st.columns([3, 1])
with h1:
    st.markdown('<div class="dash-header">SAN ANTONIO <span>MULTIFAMILY</span> INTELLIGENCE</div>', unsafe_allow_html=True)
    units_note = f"{len(df):,} permits" if not df.empty else "no permit data"
    st.markdown(
        f'<div class="dash-sub">Building Permits · {units_note} · estimated units from area · {datetime.now().strftime("%b %d, %Y")}</div>',
        unsafe_allow_html=True
    )
with h2:
    st.markdown("<br>", unsafe_allow_html=True)
    yr = st.selectbox("", ["All Time", "Last 5 Years", "Last 3 Years", "Last 12 Months", "Last 6 Months"], label_visibility="collapsed")

# Year / date filter
_year_cutoff_map = {
    "All Time":       2000,
    "Last 5 Years":   datetime.now().year - 5,
    "Last 3 Years":   datetime.now().year - 3,
    "Last 12 Months": None,
    "Last 6 Months":  None,
}
_date_cutoff_map = {
    "Last 12 Months": (datetime.now() - timedelta(days=365)).strftime('%Y-%m-%d'),
    "Last 6 Months":  (datetime.now() - timedelta(days=180)).strftime('%Y-%m-%d'),
}

if yr in _date_cutoff_map:
    _cutoff_date = _date_cutoff_map[yr]
    df_f = df[df["issue_date"] >= _cutoff_date] if not df.empty else df
    _cutoff_year = int(_cutoff_date[:4])
    dq_f = dq[dq["delivery_year"] >= _cutoff_year] if not dq.empty else dq
else:
    _cutoff_year = _year_cutoff_map.get(yr, 2000)
    df_f = df[df["delivery_year"] >= _cutoff_year] if not df.empty else df
    dq_f = dq[dq["delivery_year"] >= _cutoff_year] if not dq.empty else dq

# KPIs
k1, k2, k3, k4, k5 = st.columns(5)
for col, label, val in [
    (k1, "Est. Units Permitted",  f"{int(df_f['total_units'].sum()):,}" if not df_f.empty else "—"),
    (k2, "Projects",              f"{len(df_f):,}" if not df_f.empty else "—"),
    (k3, "Avg Est. Units",        f"{int(df_f['total_units'].mean())}" if not df_f.empty else "—"),
    (k4, "Active Submarkets",     f"{df_f['submarket_name'].nunique()}" if not df_f.empty else "—"),
    (k5, "Sell Signal Mkts",      f"{len(dc[dc['signal']=='SELL'])}"),
]:
    with col:
        st.markdown(f'<div class="kpi-card"><div class="kpi-label">{label}</div><div class="kpi-value">{val}</div></div>', unsafe_allow_html=True)

st.markdown(
    f'<div style="font-family:\'DM Mono\',monospace;font-size:0.6rem;color:{MUTED};margin-bottom:1rem;">'
    f'Note: SA permits have no unit count field. Units estimated from building area (sq ft) / 900 avg unit size, minimum 5 units.'
    f'</div>',
    unsafe_allow_html=True
)

# ─────────────────────────────────────────────────────────────────────────────
# TABS
# ─────────────────────────────────────────────────────────────────────────────
t1, t2, t3, t4, t5, t6 = st.tabs(["  MARKET OVERVIEW  ", "  SUPPLY PIPELINE  ", "  ABSORPTION  ", "  TIMING INTELLIGENCE  ", "  PERMIT BROWSER  ", "  MAP  "])

# ══════════════ TAB 1 — MARKET OVERVIEW ══════════════
with t1:
    c1, c2 = st.columns([3, 2])
    with c1:
        st.markdown('<div class="section-title">Estimated Units Permitted by Submarket</div>', unsafe_allow_html=True)
        if not df_f.empty:
            sub = df_f.groupby("submarket_name")["total_units"].sum().sort_values().reset_index()
            fig = go.Figure(go.Bar(
                x=sub["total_units"], y=sub["submarket_name"], orientation="h",
                marker=dict(
                    color=sub["total_units"],
                    colorscale=[[0, "#E5E7EB"], [0.5, "#9CA3AF"], [1, NAVY]],
                    showscale=False
                ),
                text=sub["total_units"].apply(lambda x: f"{x:,}"), textposition="outside",
                textfont=dict(size=10, color=MUTED, family="DM Mono"),
            ))
            fig.update_layout(**PLOTLY_LAYOUT, height=500, xaxis_showgrid=False, xaxis_showticklabels=False, xaxis_zeroline=False)
            st.plotly_chart(fig, use_container_width=True)
        else:
            st.info("No permit data loaded. Run: python pipeline_sanantonio.py backfill")

    with c2:
        st.markdown('<div class="section-title">Supply Pressure Score</div>', unsafe_allow_html=True)
        for _, r in dc.sort_values("score", ascending=False).iterrows():
            sc = sig_color(r["score"])
            st.markdown(f"""
            <div style="display:flex;align-items:center;gap:10px;margin-bottom:6px;padding:8px 10px;background:{CARD_BG};border:1px solid {BORDER};border-radius:4px;">
                <div style="flex:1;font-size:0.78rem;color:{TEXT};">{r['submarket_name']}</div>
                <div style="width:80px;height:4px;background:{BORDER};border-radius:2px;overflow:hidden;">
                    <div style="width:{int(r['score'])}%;height:100%;background:{sc};border-radius:2px;"></div>
                </div>
                <div style="width:28px;font-family:'DM Mono',monospace;font-size:0.7rem;color:{MUTED};text-align:right;">{r['score']:.0f}</div>
                <div style="padding:2px 6px;font-family:'DM Mono',monospace;font-size:0.62rem;border:1px solid {sc};color:{sc};border-radius:3px;">{r['signal']}</div>
            </div>
            """, unsafe_allow_html=True)

    st.markdown("<br>", unsafe_allow_html=True)
    st.markdown('<div class="section-title">Quarterly Permit Activity — All Submarkets</div>', unsafe_allow_html=True)
    if not dq_f.empty:
        qa = dq_f.groupby("delivery_yyyyq")["total_units_delivered"].sum().reset_index().sort_values("delivery_yyyyq")
        qa["rolling"] = qa["total_units_delivered"].rolling(4, min_periods=1).mean()
        fig2 = go.Figure()
        fig2.add_trace(go.Bar(x=qa["delivery_yyyyq"], y=qa["total_units_delivered"], marker_color=ACCENT, opacity=0.4, name="Quarterly"))
        fig2.add_trace(go.Scatter(x=qa["delivery_yyyyq"], y=qa["rolling"], mode="lines", line=dict(color=NAVY, width=2), name="4Q Avg"))
        fig2.update_layout(**PLOTLY_LAYOUT, height=260)
        st.plotly_chart(fig2, use_container_width=True)

# ══════════════ TAB 2 — SUPPLY PIPELINE ══════════════
with t2:
    st.markdown('<div class="section-title">Under Construction vs Historical Permit Pace</div>', unsafe_allow_html=True)
    ca, cb = st.columns(2)

    # Build pipeline table — use permit data for pace if available, else use CoStar only
    if not df_f.empty:
        pace = (
            df_f[df_f["delivery_year"] >= 2018]
            .groupby("submarket_name")["total_units"].sum()
            .div(24)
            .reset_index()
            .rename(columns={"total_units": "avg_qtr"})
        )
        pipe = dc[["submarket_name", "under_constr", "delivered_12mo", "inventory"]].copy()
        pipe = pipe.merge(pace, on="submarket_name", how="left")
        pipe["avg_qtr"] = pipe["avg_qtr"].fillna(30)
    else:
        pipe = dc[["submarket_name", "under_constr", "delivered_12mo", "inventory"]].copy()
        pipe["avg_qtr"] = 30

    pipe["months_to_deliver"] = (pipe["under_constr"] / (pipe["avg_qtr"] / 3)).clip(0, 48).round(1)
    pipe = pipe[pipe["under_constr"] > 0].sort_values("under_constr", ascending=False)

    with ca:
        fig3 = go.Figure()
        fig3.add_trace(go.Bar(x=pipe["submarket_name"], y=pipe["under_constr"], name="Under Construction", marker_color=NAVY, opacity=0.85))
        fig3.add_trace(go.Bar(x=pipe["submarket_name"], y=pipe["delivered_12mo"], name="Delivered Last 12mo", marker_color=ACCENT, opacity=0.7))
        fig3.update_layout(**PLOTLY_LAYOUT, barmode="group", height=380, xaxis_tickangle=-45)
        st.plotly_chart(fig3, use_container_width=True)

    with cb:
        st.markdown('<div class="section-title">Projected Delivery Timeline</div>', unsafe_allow_html=True)
        st.markdown(f'<div style="font-family:\'DM Mono\',monospace;font-size:0.62rem;color:{MUTED};margin-bottom:1rem;">Based on CoStar UC data and historical permit pace</div>', unsafe_allow_html=True)
        for _, r in pipe.head(13).iterrows():
            m = r["months_to_deliver"]
            uc = r["under_constr"]
            urgency_c = RED if m <= 6 else (AMBER if m <= 12 else MUTED)
            urgency_l = "IMMINENT" if m <= 6 else (f"~{m:.0f} MO")
            st.markdown(f"""
            <div style="display:flex;align-items:center;gap:10px;margin-bottom:5px;padding:8px 10px;background:{CARD_BG};border:1px solid {BORDER};border-left:2px solid {urgency_c};border-radius:4px;">
                <div style="flex:1;font-size:0.78rem;color:{TEXT};">{r['submarket_name']}</div>
                <div style="font-family:'DM Mono',monospace;font-size:0.68rem;color:{MUTED};">{uc:,.0f} UC</div>
                <div style="font-family:'DM Mono',monospace;font-size:0.68rem;color:{urgency_c};width:75px;text-align:right;">{urgency_l}</div>
            </div>
            """, unsafe_allow_html=True)

    st.markdown("<br>", unsafe_allow_html=True)
    st.markdown('<div class="section-title">Annual Permit Activity — Top 8 Submarkets</div>', unsafe_allow_html=True)
    if not df_f.empty:
        top8 = df_f.groupby("submarket_name")["total_units"].sum().nlargest(8).index.tolist()
        ann = df_f[df_f["submarket_name"].isin(top8)].groupby(["delivery_year", "submarket_name"])["total_units"].sum().reset_index()
        colors8 = [NAVY, ACCENT, "#374151", "#6B7280", "#9CA3AF", "#1A1A2E", "#C8102E", "#D97706"]
        fig4 = go.Figure()
        for i, s in enumerate(top8):
            d = ann[ann["submarket_name"] == s]
            fig4.add_trace(go.Scatter(
                x=d["delivery_year"], y=d["total_units"], name=s,
                mode="lines+markers",
                line=dict(color=colors8[i % len(colors8)], width=2),
                marker=dict(size=5)
            ))
        fig4.update_layout(**PLOTLY_LAYOUT, height=300)
        st.plotly_chart(fig4, use_container_width=True)

# ══════════════ TAB 3 — ABSORPTION ══════════════
with t3:
    st.markdown('<div class="section-title">Permit Volume vs Current Vacancy — Bubble = Units Under Construction</div>', unsafe_allow_html=True)
    if not df_f.empty:
        sub_s = df_f.groupby("submarket_name")["total_units"].sum().reset_index()
        abs_df = sub_s.merge(dc, on="submarket_name", how="inner")
        abs_df["score"] = abs_df.apply(pressure_score, axis=1)
        x_col, x_title = "total_units", "Estimated Units Permitted (historical)"
    else:
        abs_df = dc.copy()
        abs_df["score"] = abs_df.apply(pressure_score, axis=1)
        abs_df["total_units"] = abs_df["delivered_12mo"]
        x_col, x_title = "total_units", "Units Delivered (12 months, CoStar)"

    fig5 = go.Figure(go.Scatter(
        x=abs_df[x_col], y=abs_df["vacancy"] * 100,
        mode="markers+text",
        marker=dict(
            size=abs_df["under_constr"].apply(lambda x: max(8, min(x / 30, 40))),
            color=abs_df["score"],
            colorscale=[[0, GREEN], [0.5, AMBER], [1, RED]],
            showscale=True,
            colorbar=dict(title="Pressure", tickfont=dict(size=9, color=MUTED)),
            line=dict(width=1, color=BORDER)
        ),
        text=abs_df["submarket_name"],
        textposition="top center",
        textfont=dict(size=9, color=MUTED, family="DM Mono"),
        hovertemplate="<b>%{text}</b><br>Units: %{x:,}<br>Vacancy: %{y:.1f}%<extra></extra>",
    ))
    fig5.add_hline(y=10, line=dict(color=GREEN, width=1, dash="dot"), annotation_text="10% baseline", annotation_font_color=MUTED)
    fig5.add_hline(y=15, line=dict(color=AMBER, width=1, dash="dot"), annotation_text="15% caution", annotation_font_color=MUTED)
    fig5.add_hline(y=20, line=dict(color=RED, width=1, dash="dot"), annotation_text="20% oversupplied", annotation_font_color=MUTED)
    fig5.update_layout(**PLOTLY_LAYOUT, height=460, xaxis_title=x_title, yaxis_title="Vacancy Rate (%)")
    st.plotly_chart(fig5, use_container_width=True)

    st.markdown("<br>", unsafe_allow_html=True)
    st.markdown('<div class="section-title">Vacancy vs Rent Growth — Quadrant Analysis</div>', unsafe_allow_html=True)
    fig6 = go.Figure(go.Scatter(
        x=dc["vacancy"] * 100, y=dc["rent_growth"] * 100,
        mode="markers+text",
        marker=dict(
            size=12,
            color=dc["score"],
            colorscale=[[0, GREEN], [0.5, AMBER], [1, RED]],
            showscale=False,
            line=dict(width=1, color=BORDER)
        ),
        text=dc["submarket_name"],
        textposition="top center",
        textfont=dict(size=9, color=MUTED, family="DM Mono"),
        hovertemplate="<b>%{text}</b><br>Vacancy: %{x:.1f}%<br>Rent Growth: %{y:.1f}%<extra></extra>",
    ))
    fig6.add_vline(x=12, line=dict(color=BORDER, width=1, dash="dot"))
    fig6.add_hline(y=0, line=dict(color=BORDER, width=1, dash="dot"))
    for x, y, label, c in [(6, 2.5, "BUY ZONE", GREEN), (18, 2.5, "RECOVERING", AMBER), (6, -3, "WATCH", AMBER), (18, -3, "SELL ZONE", RED)]:
        fig6.add_annotation(x=x, y=y, text=label, showarrow=False, font=dict(size=8, color=c, family="DM Mono"), bgcolor="rgba(248,249,250,0.85)")
    fig6.update_layout(**PLOTLY_LAYOUT, height=380, xaxis_title="Vacancy Rate (%)", yaxis_title="Rent Growth (%)")
    st.plotly_chart(fig6, use_container_width=True)

# ══════════════ TAB 4 — TIMING INTELLIGENCE ══════════════
with t4:
    st.markdown('<div class="section-title">Buy / Hold / Sell Signal by Submarket</div>', unsafe_allow_html=True)
    st.markdown(
        f'<div style="font-family:\'DM Mono\',monospace;font-size:0.62rem;color:{MUTED};margin-bottom:1.5rem;">'
        f'Composite: vacancy (25) · deliveries (20) · pipeline (20) · rent growth (15) · absorption (10) · days on market (5) · concessions (5)'
        f'</div>',
        unsafe_allow_html=True
    )

    cs, ch, cb2 = st.columns(3)
    for col, sig_label, sc in [(cs, "SELL", RED), (ch, "HOLD", AMBER), (cb2, "BUY", GREEN)]:
        with col:
            filtered = dc[dc["signal"] == sig_label].sort_values("score", ascending=sig_label != "SELL")
            st.markdown(f'<div style="font-family:\'Inter\',sans-serif;font-size:1.2rem;font-weight:700;color:{sc};letter-spacing:0.05em;margin-bottom:1rem;padding-bottom:0.5rem;border-bottom:2px solid {sc};">{sig_label} — {len(filtered)}</div>', unsafe_allow_html=True)
            for _, r in filtered.iterrows():
                st.markdown(f"""
                <div style="padding:10px 12px;background:{CARD_BG};border:1px solid {BORDER};border-left:3px solid {sc};margin-bottom:6px;border-radius:4px;">
                    <div style="font-size:0.82rem;font-weight:600;color:{NAVY};margin-bottom:5px;">{r['submarket_name']}</div>
                    <div style="display:grid;grid-template-columns:1fr 1fr;gap:3px;">
                        <div style="font-family:'DM Mono',monospace;font-size:0.62rem;color:{MUTED};">VACANCY</div>
                        <div style="font-family:'DM Mono',monospace;font-size:0.62rem;color:{TEXT};">{r['vacancy']*100:.1f}%</div>
                        <div style="font-family:'DM Mono',monospace;font-size:0.62rem;color:{MUTED};">RENT GROWTH</div>
                        <div style="font-family:'DM Mono',monospace;font-size:0.62rem;color:{RED if r['rent_growth']<0 else GREEN};">{r['rent_growth']*100:+.1f}%</div>
                        <div style="font-family:'DM Mono',monospace;font-size:0.62rem;color:{MUTED};">UNDER CONSTR</div>
                        <div style="font-family:'DM Mono',monospace;font-size:0.62rem;color:{TEXT};">{r['under_constr']:,}</div>
                        <div style="font-family:'DM Mono',monospace;font-size:0.62rem;color:{MUTED};">ABSORPTION</div>
                        <div style="font-family:'DM Mono',monospace;font-size:0.62rem;color:{TEXT};">{r.get('absorption_12mo',0):,}</div>
                        <div style="font-family:'DM Mono',monospace;font-size:0.62rem;color:{MUTED};">CONCESSIONS</div>
                        <div style="font-family:'DM Mono',monospace;font-size:0.62rem;color:{TEXT};">{r.get('concession_pct',0)*100:.1f}%</div>
                        <div style="font-family:'DM Mono',monospace;font-size:0.62rem;color:{MUTED};">SCORE</div>
                        <div style="font-family:'DM Mono',monospace;font-size:0.62rem;color:{sc};font-weight:600;">{r['score']:.0f}/100</div>
                    </div>
                </div>
                """, unsafe_allow_html=True)

    st.markdown("<br>", unsafe_allow_html=True)
    st.markdown('<div class="section-title">Pressure Score Ranked</div>', unsafe_allow_html=True)
    sd = dc.sort_values("score", ascending=True)
    fig7 = go.Figure(go.Bar(
        x=sd["score"], y=sd["submarket_name"], orientation="h",
        marker_color=[sig_color(s) for s in sd["score"]], opacity=0.85,
        text=sd["score"].apply(lambda x: f"{x:.0f}"), textposition="outside",
        textfont=dict(size=9, color=MUTED, family="DM Mono"),
    ))
    fig7.add_vline(x=60, line=dict(color=RED, width=1, dash="dot"), annotation_text="SELL", annotation_font_color=RED)
    fig7.add_vline(x=35, line=dict(color=AMBER, width=1, dash="dot"), annotation_text="HOLD", annotation_font_color=AMBER)
    fig7.update_layout(**PLOTLY_LAYOUT, height=480, xaxis_range=[0, 110], yaxis_tickfont_size=10, yaxis_tickfont_family="DM Mono")
    st.plotly_chart(fig7, use_container_width=True)

# ══════════════ TAB 5 — PERMIT BROWSER ══════════════
with t5:
    fa, fb, fc = st.columns([2, 2, 1])
    with fa:
        search = st.text_input("", placeholder="Search address, project, or ZIP...", label_visibility="collapsed")
    with fb:
        subs = ["All Submarkets"] + sorted(df["submarket_name"].dropna().unique().tolist()) if not df.empty else ["All Submarkets"]
        sub_sel = st.selectbox("", subs, label_visibility="collapsed")
    with fc:
        min_u = st.number_input("", value=5, min_value=5, step=10, label_visibility="collapsed")

    disp = df_f.copy() if not df_f.empty else pd.DataFrame()
    if not disp.empty:
        if search:
            m = (
                disp["address"].str.contains(search, case=False, na=False) |
                disp["project_name"].str.contains(search, case=False, na=False) |
                disp["zip_code"].astype(str).str.contains(search, case=False, na=False)
            )
            disp = disp[m]
        if sub_sel != "All Submarkets":
            disp = disp[disp["submarket_name"] == sub_sel]
        if min_u > 5:
            disp = disp[disp["total_units"] >= min_u]

        st.markdown(f'<div style="font-family:\'DM Mono\',monospace;font-size:0.62rem;color:{MUTED};margin-bottom:0.5rem;">{len(disp):,} permits</div>', unsafe_allow_html=True)
        show = disp[[
            "issue_date", "address", "zip_code", "submarket_name",
            "total_units", "area_sf", "project_name", "cd", "permit_num"
        ]].rename(columns={
            "issue_date":     "Issue Date",
            "address":        "Address",
            "zip_code":       "ZIP",
            "submarket_name": "Submarket",
            "total_units":    "Est. Units",
            "area_sf":        "Area (SF)",
            "project_name":   "Project",
            "cd":             "Council Dist.",
            "permit_num":     "Permit #",
        })
        st.dataframe(show.head(500), use_container_width=True, height=500, hide_index=True)
        st.download_button(
            "Export CSV",
            disp.to_csv(index=False),
            f"sanantonio_mf_permits_{datetime.now().strftime('%Y%m%d')}.csv",
            "text/csv"
        )
    else:
        st.info("No permit data loaded. Run: python pipeline_sanantonio.py backfill")

# ══════════════ TAB 6 — MAP ══════════════
with t6:
    st.markdown('<div class="section-title">Permit Locations</div>', unsafe_allow_html=True)
    if not df_f.empty:
        map_df = df_f.dropna(subset=["latitude", "longitude"]).copy()
        map_df = map_df[(map_df["latitude"] != 0) & (map_df["longitude"] != 0)]
        if not map_df.empty:
            ma, mb = st.columns([2, 1])
            with mb:
                map_subs = ["All Submarkets"] + sorted(map_df["submarket_name"].dropna().unique().tolist())
                map_sub_sel = st.selectbox("Submarket", map_subs, label_visibility="collapsed", key="map_sub")
                map_min_units = st.slider("Minimum units", 5, 200, 5, key="map_units")
            map_show = map_df.copy()
            if map_sub_sel != "All Submarkets":
                map_show = map_show[map_show["submarket_name"] == map_sub_sel]
            if map_min_units > 5:
                map_show = map_show[map_show["total_units"] >= map_min_units]
            with mb:
                st.markdown(f'<div style="font-family:\'DM Mono\',monospace;font-size:0.62rem;color:{MUTED};margin-top:0.5rem;">{len(map_show):,} permits mapped</div>', unsafe_allow_html=True)
            with ma:
                st.map(map_show, latitude="latitude", longitude="longitude", size="total_units")
        else:
            st.info("No geocoded permits available.")
    else:
        st.info("No permit data loaded. Run: python pipeline_sanantonio.py backfill")

# ─────────────────────────────────────────────────────────────────────────────
# POWERPOINT EXPORT
# ─────────────────────────────────────────────────────────────────────────────
PPTX_NAVY  = RGBColor(0x1A, 0x1A, 0x2E)
PPTX_RED   = RGBColor(0xC8, 0x10, 0x2E)
PPTX_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PPTX_GRAY  = RGBColor(0x6B, 0x72, 0x80)

def _add_text(slide, left, top, width, height, text, font_size=12, color=None, bold=False, alignment=PP_ALIGN.LEFT):
    if color is None:
        color = PPTX_NAVY
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf

def build_pptx(dc_df, df_filtered):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PPTX_NAVY
    _add_text(slide, 0.8, 1.5, 11, 1.5, "SAN ANTONIO MULTIFAMILY INTELLIGENCE", 40, PPTX_WHITE, True, PP_ALIGN.LEFT)
    _add_text(slide, 0.8, 3.2, 8, 0.8, "Market Analysis & Investment Signals", 22, PPTX_RED, False, PP_ALIGN.LEFT)
    _add_text(slide, 0.8, 5.0, 8, 0.6, f"Matthews Real Estate Investment Services  |  {datetime.now().strftime('%B %d, %Y')}", 14, PPTX_GRAY, False, PP_ALIGN.LEFT)

    # --- Slide 2: Market KPIs ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(slide, 0.8, 0.4, 10, 0.6, "MARKET KPIs", 28, PPTX_NAVY, True)
    _add_text(slide, 0.8, 1.0, 10, 0.3, "Building Permits — Estimated Units from Area (SF) / 900", 12, PPTX_GRAY)
    kpis = [
        ("Est. Units Permitted", f"{int(df_filtered['total_units'].sum()):,}" if not df_filtered.empty else "N/A"),
        ("Projects",             f"{len(df_filtered):,}" if not df_filtered.empty else "N/A"),
        ("Avg Est. Units",       f"{int(df_filtered['total_units'].mean()):,}" if not df_filtered.empty else "N/A"),
        ("Sell Signal Mkts",     f"{len(dc_df[dc_df['signal']=='SELL'])}"),
    ]
    for i, (label, val) in enumerate(kpis):
        x = 0.8 + i * 3.0
        _add_text(slide, x, 2.0, 2.8, 0.4, label, 12, PPTX_GRAY)
        _add_text(slide, x, 2.5, 2.8, 0.8, val, 36, PPTX_NAVY, True)

    # --- Slide 3: Top SELL submarkets ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(slide, 0.8, 0.4, 10, 0.6, "SELL SIGNAL SUBMARKETS", 28, PPTX_RED, True)
    sells = dc_df[dc_df["signal"] == "SELL"].sort_values("score", ascending=False).head(5)
    headers = ["Submarket", "Score", "Vacancy", "Rent Growth", "Under Constr", "Absorption"]
    for j, h in enumerate(headers):
        _add_text(slide, 0.8 + j * 2.0, 1.5, 1.9, 0.4, h, 11, PPTX_GRAY, True)
    for i, (_, r) in enumerate(sells.iterrows()):
        y = 2.0 + i * 0.6
        vals = [
            r["submarket_name"], f"{r['score']:.0f}/100",
            f"{r['vacancy']*100:.1f}%", f"{r['rent_growth']*100:+.1f}%",
            f"{r['under_constr']:,.0f}", f"{r.get('absorption_12mo',0):,.0f}"
        ]
        for j, v in enumerate(vals):
            color = PPTX_RED if j == 1 else PPTX_NAVY
            _add_text(slide, 0.8 + j * 2.0, y, 1.9, 0.5, v, 13, color, j == 0)

    # --- Slide 4: Supply Pipeline ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(slide, 0.8, 0.4, 10, 0.6, "SUPPLY PIPELINE", 28, PPTX_NAVY, True)
    _add_text(slide, 0.8, 1.0, 10, 0.3, "Top submarkets by units under construction (CoStar)", 12, PPTX_GRAY)
    top_pipe = dc_df[dc_df["under_constr"] > 0].sort_values("under_constr", ascending=False).head(8)
    headers = ["Submarket", "Under Constr", "Delivered 12mo", "Inventory", "Vacancy"]
    for j, h in enumerate(headers):
        _add_text(slide, 0.8 + j * 2.4, 1.6, 2.3, 0.4, h, 11, PPTX_GRAY, True)
    for i, (_, r) in enumerate(top_pipe.iterrows()):
        y = 2.1 + i * 0.55
        vals = [
            r["submarket_name"], f"{r['under_constr']:,.0f}",
            f"{r['delivered_12mo']:,.0f}", f"{r['inventory']:,.0f}",
            f"{r['vacancy']*100:.1f}%"
        ]
        for j, v in enumerate(vals):
            _add_text(slide, 0.8 + j * 2.4, y, 2.3, 0.4, v, 12, PPTX_NAVY, j == 0)

    # --- Slide 5: Vacancy vs Rent Growth table ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(slide, 0.8, 0.4, 10, 0.6, "VACANCY vs RENT GROWTH", 28, PPTX_NAVY, True)
    _add_text(slide, 0.8, 1.0, 10, 0.3, "Quadrant analysis — submarkets by investment signal", 12, PPTX_GRAY)
    sorted_dc = dc_df.sort_values("score", ascending=False)
    headers = ["Submarket", "Vacancy", "Rent Growth", "Score", "Signal"]
    for j, h in enumerate(headers):
        _add_text(slide, 0.8 + j * 2.4, 1.6, 2.3, 0.4, h, 11, PPTX_GRAY, True)
    for i, (_, r) in enumerate(sorted_dc.head(13).iterrows()):
        y = 2.1 + i * 0.37
        sig_c = PPTX_RED if r["signal"] == "SELL" else (RGBColor(0xD9, 0x77, 0x06) if r["signal"] == "HOLD" else RGBColor(0x16, 0xA3, 0x4A))
        vals = [
            (r["submarket_name"], PPTX_NAVY),
            (f"{r['vacancy']*100:.1f}%", PPTX_NAVY),
            (f"{r['rent_growth']*100:+.1f}%", PPTX_RED if r["rent_growth"] < 0 else RGBColor(0x16, 0xA3, 0x4A)),
            (f"{r['score']:.0f}", PPTX_NAVY),
            (r["signal"], sig_c),
        ]
        for j, (v, c) in enumerate(vals):
            _add_text(slide, 0.8 + j * 2.4, y, 2.3, 0.3, v, 11, c, j == 0 or j == 4)

    # --- Slide 6: Methodology ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(slide, 0.8, 0.4, 10, 0.6, "METHODOLOGY & DATA SOURCES", 28, PPTX_NAVY, True)
    methodology = (
        "Supply Pressure Score (0-100)\n"
        "Composite of seven weighted factors:\n"
        "  - Vacancy Rate: 25 pts\n"
        "  - 12-Month Deliveries vs Inventory: 20 pts\n"
        "  - Under Construction vs Inventory: 20 pts\n"
        "  - Rent Growth (inverted): 15 pts\n"
        "  - Absorption vs Deliveries: 10 pts\n"
        "  - Avg Days on Market: 5 pts\n"
        "  - Concession Rate: 5 pts\n\n"
        "Signals: BUY (<35) | HOLD (35-59) | SELL (60+)\n\n"
        "Data Sources:\n"
        "  - City of San Antonio Open Data Portal (data.sanantonio.gov, CKAN API)\n"
        "  - Permit Type: Comm New Building Permit, Work Type: New\n"
        "  - CoStar Group (vacancy, rent, absorption, pipeline)\n"
        "  - Units estimated from building area (sq ft) / 900 avg unit size, min 5 units\n"
        "  - Deduplicated by address + unit count to avoid duplicate entries"
    )
    _add_text(slide, 0.8, 1.2, 11, 5.0, methodology, 14, PPTX_NAVY)
    _add_text(slide, 0.8, 6.5, 11, 0.5, "Matthews Real Estate Investment Services  |  Confidential", 11, PPTX_GRAY)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf

with st.sidebar:
    st.markdown(f'<div style="font-family:\'DM Mono\',monospace;font-size:0.65rem;color:{MUTED};letter-spacing:0.15em;text-transform:uppercase;margin-top:1.5rem;margin-bottom:0.5rem;">Export</div>', unsafe_allow_html=True)
    pptx_buf = build_pptx(dc, df_f)
    st.download_button(
        label="Export to PowerPoint",
        data=pptx_buf,
        file_name=f"sanantonio_mf_intelligence_{datetime.now().strftime('%Y%m%d')}.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

# ─────────────────────────────────────────────────────────────────────────────
# FOOTER
# ─────────────────────────────────────────────────────────────────────────────
st.markdown(f"""
<div style="margin-top:3rem;padding-top:1rem;border-top:1px solid {BORDER};display:flex;justify-content:space-between;align-items:center;">
    <div style="font-family:'DM Mono',monospace;font-size:0.6rem;color:{MUTED};">DATA: data.sanantonio.gov (CKAN) · CoStar Group · Comm New Building Permits · Units estimated from area</div>
    <div style="font-family:'DM Mono',monospace;font-size:0.6rem;color:{MUTED};letter-spacing:0.12em;">MATTHEWS REAL ESTATE INVESTMENT SERVICES</div>
</div>
""", unsafe_allow_html=True)
