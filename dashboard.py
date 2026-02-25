"""
Austin Multifamily Intelligence Dashboard
==========================================
Run with: streamlit run dashboard.py

Tabs:
  1. Market Overview    — Supply pressure score by submarket
  2. Supply Pipeline    — Under construction + projected delivery dates
  3. Absorption         — Delivery volume vs vacancy over time
  4. Timing Intelligence — Buy/sell signals by submarket
  5. Permit Browser     — Raw CO data, searchable
"""

import io
import json
import os
from datetime import datetime, timedelta
from typing import Optional

import pandas as pd
import plotly.graph_objects as go
import streamlit as st
import psycopg2
import psycopg2.extras
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

load_dotenv()

# Streamlit Cloud secrets → env var → local default
try:
    DB_DSN = st.secrets["DATABASE_URL"]
except (KeyError, FileNotFoundError, AttributeError):
    DB_DSN = os.getenv("DATABASE_URL", "postgresql://postgres:postgres@localhost:5432/austin_co")

# ─────────────────────────────────────────────────────────────────────────────
# COSTAR DATA — update these from your GeographyList.xlsx periodically
# ─────────────────────────────────────────────────────────────────────────────
COSTAR_DATA = {
    "North Austin": {
        "vacancy": 0.140,
        "rent_growth": -0.063,
        "inventory": 31903,
        "under_constr": 1065,
        "delivered_12mo": 374,
        "asking_rent": 1395,
        "absorption_12mo": 1379,
        "avg_days_on_market": 58,
        "concession_pct": 0.065,
    },
    "Northwest Austin": {
        "vacancy": 0.110,
        "rent_growth": -0.059,
        "inventory": 35400,
        "under_constr": 586,
        "delivered_12mo": 129,
        "asking_rent": 1343,
        "absorption_12mo": 364,
        "avg_days_on_market": 48,
        "concession_pct": 0.055,
    },
    "East Austin": {
        "vacancy": 0.138,
        "rent_growth": -0.041,
        "inventory": 19201,
        "under_constr": 2328,
        "delivered_12mo": 792,
        "asking_rent": 1774,
        "absorption_12mo": 1636,
        "avg_days_on_market": 55,
        "concession_pct": 0.050,
    },
    "Downtown Austin": {
        "vacancy": 0.138,
        "rent_growth": 0.012,
        "inventory": 9001,
        "under_constr": 352,
        "delivered_12mo": 1254,
        "asking_rent": 3575,
        "absorption_12mo": 994,
        "avg_days_on_market": 62,
        "concession_pct": 0.045,
    },
    "Pflugerville": {
        "vacancy": 0.139,
        "rent_growth": -0.067,
        "inventory": 23825,
        "under_constr": 444,
        "delivered_12mo": 752,
        "asking_rent": 1386,
        "absorption_12mo": 1074,
        "avg_days_on_market": 60,
        "concession_pct": 0.070,
    },
    "Northeast Austin": {
        "vacancy": 0.213,
        "rent_growth": -0.057,
        "inventory": 19387,
        "under_constr": 2436,
        "delivered_12mo": 3102,
        "asking_rent": 1397,
        "absorption_12mo": 2138,
        "avg_days_on_market": 82,
        "concession_pct": 0.105,
    },
    "South Austin": {
        "vacancy": 0.119,
        "rent_growth": -0.056,
        "inventory": 21649,
        "under_constr": 1063,
        "delivered_12mo": 930,
        "asking_rent": 1386,
        "absorption_12mo": 955,
        "avg_days_on_market": 52,
        "concession_pct": 0.060,
    },
    "Round Rock": {
        "vacancy": 0.105,
        "rent_growth": -0.064,
        "inventory": 21390,
        "under_constr": 217,
        "delivered_12mo": 252,
        "asking_rent": 1426,
        "absorption_12mo": 1400,
        "avg_days_on_market": 44,
        "concession_pct": 0.055,
    },
    "Midtown Austin": {
        "vacancy": 0.127,
        "rent_growth": -0.021,
        "inventory": 17027,
        "under_constr": 1681,
        "delivered_12mo": 426,
        "asking_rent": 1562,
        "absorption_12mo": 808,
        "avg_days_on_market": 54,
        "concession_pct": 0.040,
    },
    "Georgetown-Leander": {
        "vacancy": 0.166,
        "rent_growth": -0.067,
        "inventory": 17489,
        "under_constr": 574,
        "delivered_12mo": 1251,
        "asking_rent": 1514,
        "absorption_12mo": 2053,
        "avg_days_on_market": 65,
        "concession_pct": 0.075,
    },
    "Southeast Austin": {
        "vacancy": 0.209,
        "rent_growth": -0.052,
        "inventory": 16442,
        "under_constr": 571,
        "delivered_12mo": 2335,
        "asking_rent": 1397,
        "absorption_12mo": 1416,
        "avg_days_on_market": 78,
        "concession_pct": 0.095,
    },
    "Riverside": {
        "vacancy": 0.116,
        "rent_growth": -0.065,
        "inventory": 18784,
        "under_constr": 298,
        "delivered_12mo": 401,
        "asking_rent": 1377,
        "absorption_12mo": 542,
        "avg_days_on_market": 55,
        "concession_pct": 0.068,
    },
    "Southwest Austin": {
        "vacancy": 0.113,
        "rent_growth": -0.034,
        "inventory": 13882,
        "under_constr": 949,
        "delivered_12mo": 803,
        "asking_rent": 1687,
        "absorption_12mo": 238,
        "avg_days_on_market": 58,
        "concession_pct": 0.042,
    },
    "Cedar Park": {
        "vacancy": 0.116,
        "rent_growth": -0.054,
        "inventory": 15829,
        "under_constr": 0,
        "delivered_12mo": 391,
        "asking_rent": 1438,
        "absorption_12mo": 905,
        "avg_days_on_market": 50,
        "concession_pct": 0.058,
    },
    "South Central Austin": {
        "vacancy": 0.116,
        "rent_growth": -0.035,
        "inventory": 13679,
        "under_constr": 572,
        "delivered_12mo": 570,
        "asking_rent": 1740,
        "absorption_12mo": 210,
        "avg_days_on_market": 53,
        "concession_pct": 0.040,
    },
    "Buda-Kyle": {
        "vacancy": 0.157,
        "rent_growth": -0.046,
        "inventory": 11355,
        "under_constr": 280,
        "delivered_12mo": 798,
        "asking_rent": 1445,
        "absorption_12mo": 992,
        "avg_days_on_market": 63,
        "concession_pct": 0.072,
    },
    "San Marcos": {
        "vacancy": 0.214,
        "rent_growth": -0.051,
        "inventory": 10828,
        "under_constr": 807,
        "delivered_12mo": 725,
        "asking_rent": 1261,
        "absorption_12mo": 1151,
        "avg_days_on_market": 80,
        "concession_pct": 0.090,
    },
    "Far North Austin": {
        "vacancy": 0.262,
        "rent_growth": -0.033,
        "inventory": 4032,
        "under_constr": 336,
        "delivered_12mo": 917,
        "asking_rent": 1561,
        "absorption_12mo": 701,
        "avg_days_on_market": 95,
        "concession_pct": 0.115,
    },
    "Lake Travis": {
        "vacancy": 0.141,
        "rent_growth": -0.030,
        "inventory": 3978,
        "under_constr": 0,
        "delivered_12mo": 322,
        "asking_rent": 1813,
        "absorption_12mo": 38,
        "avg_days_on_market": 68,
        "concession_pct": 0.048,
    },
    "Central Austin": {
        "vacancy": 0.092,
        "rent_growth": -0.021,
        "inventory": 3971,
        "under_constr": 0,
        "delivered_12mo": 38,
        "asking_rent": 1578,
        "absorption_12mo": -10,
        "avg_days_on_market": 42,
        "concession_pct": 0.030,
    },
    "West Austin": {
        "vacancy": 0.061,
        "rent_growth": -0.006,
        "inventory": 2152,
        "under_constr": 168,
        "delivered_12mo": 0,
        "asking_rent": 2009,
        "absorption_12mo": 14,
        "avg_days_on_market": 32,
        "concession_pct": 0.018,
    },
    "Far West Austin": {
        "vacancy": 0.054,
        "rent_growth": 0.008,
        "inventory": 149,
        "under_constr": 0,
        "delivered_12mo": 0,
        "asking_rent": 1645,
        "absorption_12mo": 0,
        "avg_days_on_market": 28,
        "concession_pct": 0.015,
    },
}

# ─────────────────────────────────────────────────────────────────────────────
# MATTHEWS BRAND COLORS
# ─────────────────────────────────────────────────────────────────────────────
ACCENT    = "#C8102E"   # Matthews red — highlights, active tabs, key metrics
NAVY      = "#1A1A2E"   # Deep navy/charcoal — headers and primary text
GREEN     = "#16A34A"   # Success / BUY
AMBER     = "#D97706"   # Warning / HOLD
RED       = "#C8102E"   # Danger / SELL (same as accent)
BG        = "#FFFFFF"   # Primary background — white
CARD_BG   = "#F8F9FA"   # Card background — light gray
BORDER    = "#E5E7EB"   # Subtle gray border
TEXT      = "#1A1A2E"   # Primary text — navy
MUTED     = "#6B7280"   # Secondary text — gray

PLOTLY_LAYOUT = dict(
    paper_bgcolor=BG,
    plot_bgcolor=BG,
    font=dict(color=TEXT, family="'DM Sans', sans-serif"),
    xaxis=dict(gridcolor=BORDER, zerolinecolor=BORDER, tickfont=dict(color=MUTED)),
    yaxis=dict(gridcolor=BORDER, zerolinecolor=BORDER, tickfont=dict(color=MUTED)),
    legend=dict(bgcolor="rgba(255,255,255,0.9)", bordercolor=BORDER, font=dict(color=TEXT)),
    margin=dict(t=40, r=20, b=40, l=60),
)

# ─────────────────────────────────────────────────────────────────────────────
# PAGE CONFIG
# ─────────────────────────────────────────────────────────────────────────────
st.set_page_config(page_title="Austin MF Intelligence", page_icon="🏢", layout="wide")

st.markdown(f"""
<style>
@import url('https://fonts.googleapis.com/css2?family=DM+Mono:wght@300;400;500&family=Inter:wght@300;400;500;600;700&family=DM+Sans:wght@300;400;500;600&display=swap');

html, body, [class*="css"] {{ background-color: {BG}; color: {TEXT}; font-family: 'DM Sans', sans-serif; }}
.main {{ background-color: {BG}; }}
.stApp {{ margin-top: 0; padding-top: 0; }}
.block-container {{ padding-top: 0 !important; padding-right: 1rem; padding-bottom: 1.5rem; padding-left: 1rem; max-width: 1600px; }}
header[data-testid="stHeader"] {{ display: none !important; }}
.stDeployButton {{ display: none !important; }}
#MainMenu {{ display: none !important; }}
div[data-testid="stToolbar"] {{ display: none !important; }}
div[data-testid="stDecoration"] {{ display: none !important; }}
.uploadedFile {{ display: none !important; }}
section[data-testid="stFileUploadDropzone"] {{ display: none !important; }}
div[data-testid="stFileDropzoneInput"] {{ display: none !important; }}
[data-testid="stFileUploader"] {{ display: none !important; }}
.stFileDropzone {{ display: none !important; }}
div[data-testid="stAppViewBlockContainer"] [data-testid="stFileUploader"] {{ display: none !important; }}
.drag-drop-container {{ display: none !important; }}
[data-testid="stFileDropzone"] {{ display: none !important; }}
.dash-header {{ font-family: 'Inter', sans-serif; font-size: clamp(1.2rem, 3vw, 2.2rem); font-weight: 700; letter-spacing: 0.02em; color: {NAVY}; line-height: 1.1; white-space: nowrap; overflow: hidden; text-overflow: ellipsis; }}
.dash-header span {{ color: {ACCENT}; }}
.dash-sub {{ font-family: 'DM Mono', monospace; font-size: 0.72rem; color: {MUTED}; letter-spacing: 0.18em; text-transform: uppercase; margin-bottom: 1.5rem; }}
.kpi-card {{ background: {CARD_BG}; border: 1px solid {BORDER}; border-left: 3px solid {ACCENT}; border-radius: 4px; padding: 1rem 1.2rem; margin-bottom: 1rem; }}
.kpi-label {{ font-family: 'DM Mono', monospace; font-size: 0.65rem; color: {MUTED}; letter-spacing: 0.15em; text-transform: uppercase; margin-bottom: 0.4rem; }}
.kpi-value {{ font-family: 'Inter', sans-serif; font-size: 1.9rem; font-weight: 700; color: {NAVY}; line-height: 1; }}
.section-title {{ font-family: 'DM Mono', monospace; font-size: 0.65rem; color: {MUTED}; letter-spacing: 0.2em; text-transform: uppercase; border-bottom: 1px solid {BORDER}; padding-bottom: 0.5rem; margin-bottom: 1rem; }}
.stTabs [data-baseweb="tab-list"] {{ background-color: {BG}; border-bottom: 1px solid {BORDER}; gap: 0; }}
.stTabs [data-baseweb="tab"] {{ font-family: 'DM Mono', monospace; font-size: 0.7rem; letter-spacing: 0.12em; text-transform: uppercase; color: {MUTED}; padding: 0.75rem 1.5rem; border-bottom: 2px solid transparent; background: transparent; }}
.stTabs [aria-selected="true"] {{ color: {ACCENT} !important; border-bottom: 2px solid {ACCENT} !important; background: transparent !important; }}
.stDataFrame {{ border: 1px solid {BORDER}; border-radius: 4px; }}
input, select, textarea {{ background-color: {CARD_BG} !important; color: {TEXT} !important; border: 1px solid {BORDER} !important; border-radius: 4px !important; }}
::-webkit-scrollbar {{ width: 4px; height: 4px; }}
::-webkit-scrollbar-track {{ background: {BG}; }}
::-webkit-scrollbar-thumb {{ background: {BORDER}; border-radius: 2px; }}
</style>
""", unsafe_allow_html=True)

# ─────────────────────────────────────────────────────────────────────────────
# DATA
# ─────────────────────────────────────────────────────────────────────────────
@st.cache_data(ttl=300)
def load_permits():
    conn = psycopg2.connect(DB_DSN)
    df = pd.read_sql("""
        SELECT permit_num, masterpermitnum, permit_class, issue_date, address,
               zip_code, latitude, longitude, total_units, project_name,
               work_class, submarket_name,
               delivery_year, delivery_quarter, delivery_yyyyq
        FROM co_projects
        WHERE issue_date IS NOT NULL
        ORDER BY issue_date DESC
    """, conn)
    conn.close()
    df["issue_date"] = pd.to_datetime(df["issue_date"])
    return df

@st.cache_data(ttl=300)
def load_quarterly():
    conn = psycopg2.connect(DB_DSN)
    df = pd.read_sql("""
        SELECT submarket_name, delivery_year, delivery_quarter,
               delivery_yyyyq, project_count, total_units_delivered
        FROM submarket_deliveries ORDER BY delivery_yyyyq
    """, conn)
    conn.close()
    return df

@st.cache_data(ttl=300)
def load_submarket_boundaries():
    """Load submarket polygon boundaries as GeoJSON from PostGIS."""
    conn = psycopg2.connect(DB_DSN)
    cur = conn.cursor()
    cur.execute("""
        SELECT submarket_name, ST_AsGeoJSON(geom)::text as geojson
        FROM costar_submarkets
        WHERE geom IS NOT NULL
        ORDER BY submarket_name
    """)
    features = []
    for name, geojson_str in cur.fetchall():
        features.append({
            "type": "Feature",
            "properties": {"submarket_name": name},
            "geometry": json.loads(geojson_str)
        })
    conn.close()
    return {"type": "FeatureCollection", "features": features}

def get_costar_df():
    rows = [{"submarket_name": k, **v} for k, v in COSTAR_DATA.items()]
    return pd.DataFrame(rows)

def pressure_score(row):
    # Existing factors (scaled back to make room for new signals)
    v = min((row["vacancy"] - 0.08) / 0.15, 1.0) * 25          # 25pts vacancy
    d = min(row["delivered_12mo"] / max(row["inventory"], 1) / 0.12, 1.0) * 20  # 20pts deliveries
    u = min(row["under_constr"] / max(row["inventory"], 1) / 0.15, 1.0) * 20   # 20pts pipeline
    r = min((-row["rent_growth"]) / 0.08, 1.0) * 15             # 15pts rent growth

    # New factors
    # Absorption: low absorption vs deliveries = pressure (clamped 0-10)
    absorption_ratio = row.get("absorption_12mo", 0) / max(row["delivered_12mo"], 1)
    a = max(0.0, min((1 - absorption_ratio) / 0.5, 1.0)) * 10   # 10pts absorption

    # Days on market: above 45 days = pressure signal (clamped 0-5)
    dom = row.get("avg_days_on_market", 45)
    dom_score = max(0.0, min((dom - 45) / 60, 1.0)) * 5         # 5pts days on market

    # Concessions: above 4% = distress signal (clamped 0-5)
    conc = row.get("concession_pct", 0)
    conc_score = max(0.0, min(max(conc - 0.04, 0) / 0.10, 1.0)) * 5  # 5pts concessions

    return round(max(0, v + d + u + r + a + dom_score + conc_score), 1)

def sig(score):
    if score >= 60: return "SELL"
    if score >= 35: return "HOLD"
    return "BUY"

def sig_color(score):
    if score >= 60: return RED
    if score >= 35: return AMBER
    return GREEN

try:
    df = load_permits()
    dq = load_quarterly()
    db_ok = True
except Exception as e:
    st.error(f"DB error: {e}")
    df = pd.DataFrame()
    dq = pd.DataFrame()
    db_ok = False

dc = get_costar_df()
dc["score"] = dc.apply(pressure_score, axis=1)
dc["signal"] = dc["score"].apply(sig)

# ─────────────────────────────────────────────────────────────────────────────
# HEADER
# ─────────────────────────────────────────────────────────────────────────────
h1, h2 = st.columns([3, 1])
with h1:
    st.markdown('<div class="dash-header">AUSTIN <span>MULTIFAMILY</span> INTELLIGENCE</div>', unsafe_allow_html=True)
    st.markdown(f'<div class="dash-sub">Certificates of Occupancy · {len(df):,} permits · {datetime.now().strftime("%b %d, %Y")}</div>', unsafe_allow_html=True)
with h2:
    st.markdown("<br>", unsafe_allow_html=True)
    yr = st.selectbox("", ["All Time", "Last 5 Years", "Last 3 Years", "Last 12 Months", "Last 6 Months"], label_visibility="collapsed")

# Year-based cutoffs (for delivery_year filtering)
_year_cutoff_map = {
    "All Time":      None,
    "Last 5 Years":  datetime.now().year - 5,
    "Last 3 Years":  datetime.now().year - 3,
    "Last 12 Months": None,
    "Last 6 Months":  None,
}
# Date-based cutoffs (for issue_date filtering — more precise)
_date_cutoff_map = {
    "Last 12 Months": (datetime.now() - timedelta(days=365)).strftime('%Y-%m-%d'),
    "Last 6 Months":  (datetime.now() - timedelta(days=180)).strftime('%Y-%m-%d'),
}

if yr in _date_cutoff_map:
    _cutoff_date = _date_cutoff_map[yr]
    df_f = df[df["issue_date"] >= _cutoff_date] if not df.empty else df
    # For quarterly data, approximate using delivery_year from the date cutoff
    _cutoff_year = int(_cutoff_date[:4])
    dq_f = dq[dq["delivery_year"] >= _cutoff_year] if not dq.empty else dq
else:
    _cutoff_year = _year_cutoff_map.get(yr)
    if _cutoff_year is not None:
        df_f = df[df["delivery_year"] >= _cutoff_year] if not df.empty else df
        dq_f = dq[dq["delivery_year"] >= _cutoff_year] if not dq.empty else dq
    else:
        df_f = df.copy() if not df.empty else df
        dq_f = dq.copy() if not dq.empty else dq

# KPIs
k1, k2, k3, k4, k5 = st.columns(5)
for col, label, val in [
    (k1, "Units Delivered",   f"{int(df_f['total_units'].sum()):,}" if not df_f.empty else "—"),
    (k2, "Projects",          f"{len(df_f):,}" if not df_f.empty else "—"),
    (k3, "Avg Project Size",  f"{int(df_f['total_units'].mean())}" if not df_f.empty else "—"),
    (k4, "Active Submarkets", f"{df_f['submarket_name'].nunique()}" if not df_f.empty else "—"),
    (k5, "Sell Signal Mkts",  f"{len(dc[dc['signal']=='SELL'])}"),
]:
    with col:
        st.markdown(f'<div class="kpi-card"><div class="kpi-label">{label}</div><div class="kpi-value">{val}</div></div>', unsafe_allow_html=True)

st.markdown("<br>", unsafe_allow_html=True)

# ─────────────────────────────────────────────────────────────────────────────
# TABS
# ─────────────────────────────────────────────────────────────────────────────
t1, t2, t3, t4, t5, t6 = st.tabs(["  MARKET OVERVIEW  ", "  SUPPLY PIPELINE  ", "  ABSORPTION  ", "  TIMING INTELLIGENCE  ", "  PERMIT BROWSER  ", "  MAP  "])

# ══════════════ TAB 1 ══════════════
with t1:
    c1, c2 = st.columns([3, 2])
    with c1:
        st.markdown('<div class="section-title">Units Delivered by Submarket</div>', unsafe_allow_html=True)
        if not df_f.empty:
            sub = df_f.groupby("submarket_name")["total_units"].sum().sort_values().reset_index()
            fig = go.Figure(go.Bar(
                x=sub["total_units"], y=sub["submarket_name"], orientation="h",
                marker=dict(
                    color=sub["total_units"],
                    colorscale=[[0, "#E5E7EB"], [0.5, "#9CA3AF"], [1, NAVY]],
                    showscale=False
                ),
                text=sub["total_units"].apply(lambda x: f"{x:,}"), textposition="outside",
                textfont=dict(size=10, color=MUTED, family="DM Mono"),
            ))
            fig.update_layout(**PLOTLY_LAYOUT, height=500, xaxis_showgrid=False, xaxis_showticklabels=False, xaxis_zeroline=False)
            st.plotly_chart(fig, use_container_width=True)

    with c2:
        st.markdown('<div class="section-title">Supply Pressure Score</div>', unsafe_allow_html=True)
        for _, r in dc.sort_values("score", ascending=False).iterrows():
            sc = sig_color(r["score"])
            st.markdown(f"""
            <div style="display:flex;align-items:center;gap:10px;margin-bottom:6px;padding:8px 10px;background:{CARD_BG};border:1px solid {BORDER};border-radius:4px;">
                <div style="flex:1;font-size:0.78rem;color:{TEXT};">{r['submarket_name']}</div>
                <div style="width:80px;height:4px;background:{BORDER};border-radius:2px;overflow:hidden;">
                    <div style="width:{int(r['score'])}%;height:100%;background:{sc};border-radius:2px;"></div>
                </div>
                <div style="width:28px;font-family:'DM Mono',monospace;font-size:0.7rem;color:{MUTED};text-align:right;">{r['score']:.0f}</div>
                <div style="padding:2px 6px;font-family:'DM Mono',monospace;font-size:0.62rem;border:1px solid {sc};color:{sc};border-radius:3px;">{r['signal']}</div>
            </div>
            """, unsafe_allow_html=True)

    st.markdown("<br>", unsafe_allow_html=True)
    st.markdown('<div class="section-title">Quarterly Deliveries — All Submarkets</div>', unsafe_allow_html=True)
    if not dq_f.empty:
        qa = dq_f.groupby("delivery_yyyyq")["total_units_delivered"].sum().reset_index().sort_values("delivery_yyyyq")
        qa["rolling"] = qa["total_units_delivered"].rolling(4, min_periods=1).mean()
        fig2 = go.Figure()
        fig2.add_trace(go.Bar(x=qa["delivery_yyyyq"], y=qa["total_units_delivered"], marker_color=ACCENT, opacity=0.4, name="Quarterly"))
        fig2.add_trace(go.Scatter(x=qa["delivery_yyyyq"], y=qa["rolling"], mode="lines", line=dict(color=NAVY, width=2), name="4Q Avg"))
        fig2.update_layout(**PLOTLY_LAYOUT, height=260)
        st.plotly_chart(fig2, use_container_width=True)

# ══════════════ TAB 2 ══════════════
with t2:
    st.markdown('<div class="section-title">Under Construction vs Historical Delivery Pace</div>', unsafe_allow_html=True)
    ca, cb = st.columns(2)

    with ca:
        if not df_f.empty:
            pace = df_f[df_f["delivery_year"] >= 2018].groupby("submarket_name")["total_units"].sum().div(24).reset_index().rename(columns={"total_units": "avg_qtr"})
            pipe = dc[["submarket_name","under_constr","delivered_12mo","inventory"]].copy()
            pipe = pipe.merge(pace, on="submarket_name", how="left")
            pipe["avg_qtr"] = pipe["avg_qtr"].fillna(50)
            pipe["months_to_deliver"] = (pipe["under_constr"] / (pipe["avg_qtr"] / 3)).clip(0, 48).round(1)
            pipe = pipe[pipe["under_constr"] > 0].sort_values("under_constr", ascending=False)

            fig3 = go.Figure()
            fig3.add_trace(go.Bar(x=pipe["submarket_name"], y=pipe["under_constr"], name="Under Construction", marker_color=NAVY, opacity=0.85))
            fig3.add_trace(go.Bar(x=pipe["submarket_name"], y=pipe["delivered_12mo"], name="Delivered Last 12mo", marker_color=ACCENT, opacity=0.7))
            fig3.update_layout(**PLOTLY_LAYOUT, barmode="group", height=380, xaxis_tickangle=-45)
            st.plotly_chart(fig3, use_container_width=True)

    with cb:
        st.markdown('<div class="section-title">Projected Delivery Timeline</div>', unsafe_allow_html=True)
        st.markdown(f'<div style="font-family:\'DM Mono\',monospace;font-size:0.62rem;color:{MUTED};margin-bottom:1rem;">Based on historical CO pace — not CoStar estimates</div>', unsafe_allow_html=True)
        if not df_f.empty:
            for _, r in pipe.head(14).iterrows():
                m = r["months_to_deliver"]
                uc = r["under_constr"]
                urgency_c = RED if m <= 6 else (AMBER if m <= 12 else MUTED)
                urgency_l = "IMMINENT" if m <= 6 else (f"~{m:.0f} MO")
                st.markdown(f"""
                <div style="display:flex;align-items:center;gap:10px;margin-bottom:5px;padding:8px 10px;background:{CARD_BG};border:1px solid {BORDER};border-left:2px solid {urgency_c};border-radius:4px;">
                    <div style="flex:1;font-size:0.78rem;color:{TEXT};">{r['submarket_name']}</div>
                    <div style="font-family:'DM Mono',monospace;font-size:0.68rem;color:{MUTED};">{uc:,.0f} UC</div>
                    <div style="font-family:'DM Mono',monospace;font-size:0.68rem;color:{urgency_c};width:75px;text-align:right;">{urgency_l}</div>
                </div>
                """, unsafe_allow_html=True)

    st.markdown("<br>", unsafe_allow_html=True)
    st.markdown('<div class="section-title">Annual Delivery Volume — Top 8 Submarkets</div>', unsafe_allow_html=True)
    if not df_f.empty:
        top8 = df_f.groupby("submarket_name")["total_units"].sum().nlargest(8).index.tolist()
        ann = df_f[df_f["submarket_name"].isin(top8)].groupby(["delivery_year","submarket_name"])["total_units"].sum().reset_index()
        colors8 = [NAVY, ACCENT, "#374151", "#6B7280", "#9CA3AF", "#1A1A2E", "#C8102E", "#D97706"]
        fig4 = go.Figure()
        for i, s in enumerate(top8):
            d = ann[ann["submarket_name"] == s]
            fig4.add_trace(go.Scatter(x=d["delivery_year"], y=d["total_units"], name=s, mode="lines+markers", line=dict(color=colors8[i % len(colors8)], width=2), marker=dict(size=5)))
        fig4.update_layout(**PLOTLY_LAYOUT, height=300)
        st.plotly_chart(fig4, use_container_width=True)

# ══════════════ TAB 3 ══════════════
with t3:
    st.markdown('<div class="section-title">Delivery Volume vs Current Vacancy — Bubble = Units Under Construction</div>', unsafe_allow_html=True)
    if not df_f.empty:
        sub_s = df_f.groupby("submarket_name")["total_units"].sum().reset_index()
        abs_df = sub_s.merge(dc, on="submarket_name", how="inner")
        abs_df["score"] = abs_df.apply(pressure_score, axis=1)

        fig5 = go.Figure(go.Scatter(
            x=abs_df["total_units"], y=abs_df["vacancy"] * 100,
            mode="markers+text",
            marker=dict(size=abs_df["under_constr"].apply(lambda x: max(8, min(x/50,40))),
                       color=abs_df["score"], colorscale=[[0,GREEN],[0.5,AMBER],[1,RED]],
                       showscale=True, colorbar=dict(title="Pressure", tickfont=dict(size=9,color=MUTED)),
                       line=dict(width=1,color=BORDER)),
            text=abs_df["submarket_name"].apply(lambda x: x.replace(" Austin","").replace(" County","")),
            textposition="top center", textfont=dict(size=9,color=MUTED,family="DM Mono"),
            hovertemplate="<b>%{text}</b><br>Units: %{x:,}<br>Vacancy: %{y:.1f}%<extra></extra>",
        ))
        fig5.add_hline(y=10, line=dict(color=GREEN,width=1,dash="dot"), annotation_text="10% baseline", annotation_font_color=MUTED)
        fig5.add_hline(y=15, line=dict(color=AMBER,width=1,dash="dot"), annotation_text="15% caution", annotation_font_color=MUTED)
        fig5.add_hline(y=20, line=dict(color=RED,width=1,dash="dot"), annotation_text="20% oversupplied", annotation_font_color=MUTED)
        fig5.update_layout(**PLOTLY_LAYOUT, height=460, xaxis_title="Total Units Delivered (historical)", yaxis_title="Vacancy Rate (%)")
        st.plotly_chart(fig5, use_container_width=True)

    st.markdown("<br>", unsafe_allow_html=True)
    st.markdown('<div class="section-title">Vacancy vs Rent Growth — Quadrant Analysis</div>', unsafe_allow_html=True)
    fig6 = go.Figure(go.Scatter(
        x=dc["vacancy"]*100, y=dc["rent_growth"]*100,
        mode="markers+text",
        marker=dict(size=12, color=dc["score"], colorscale=[[0,GREEN],[0.5,AMBER],[1,RED]], showscale=False, line=dict(width=1,color=BORDER)),
        text=dc["submarket_name"].apply(lambda x: x.replace(" Austin","").replace(" County","")),
        textposition="top center", textfont=dict(size=9,color=MUTED,family="DM Mono"),
        hovertemplate="<b>%{text}</b><br>Vacancy: %{x:.1f}%<br>Rent Growth: %{y:.1f}%<extra></extra>",
    ))
    fig6.add_vline(x=14, line=dict(color=BORDER,width=1,dash="dot"))
    fig6.add_hline(y=0, line=dict(color=BORDER,width=1,dash="dot"))
    for x, y, label, c in [(8,3,"BUY ZONE",GREEN),(20,3,"RECOVERING",AMBER),(8,-4,"WATCH",AMBER),(20,-4,"SELL ZONE",RED)]:
        fig6.add_annotation(x=x,y=y,text=label,showarrow=False,font=dict(size=8,color=c,family="DM Mono"),bgcolor="rgba(248,249,250,0.85)")
    fig6.update_layout(**PLOTLY_LAYOUT, height=380, xaxis_title="Vacancy Rate (%)", yaxis_title="Rent Growth (%)")
    st.plotly_chart(fig6, use_container_width=True)

# ══════════════ TAB 4 ══════════════
with t4:
    st.markdown('<div class="section-title">Buy / Hold / Sell Signal by Submarket</div>', unsafe_allow_html=True)
    st.markdown(f'<div style="font-family:\'DM Mono\',monospace;font-size:0.62rem;color:{MUTED};margin-bottom:1.5rem;">Composite: vacancy (25) · deliveries (20) · pipeline (20) · rent growth (15) · absorption (10) · days on market (5) · concessions (5)</div>', unsafe_allow_html=True)

    cs, ch, cb2 = st.columns(3)
    for col, sig_label, sc in [(cs,"SELL",RED),(ch,"HOLD",AMBER),(cb2,"BUY",GREEN)]:
        with col:
            filtered = dc[dc["signal"] == sig_label].sort_values("score", ascending=sig_label!="SELL")
            st.markdown(f'<div style="font-family:\'Inter\',sans-serif;font-size:1.2rem;font-weight:700;color:{sc};letter-spacing:0.05em;margin-bottom:1rem;padding-bottom:0.5rem;border-bottom:2px solid {sc};">{sig_label} — {len(filtered)}</div>', unsafe_allow_html=True)
            for _, r in filtered.iterrows():
                st.markdown(f"""
                <div style="padding:10px 12px;background:{CARD_BG};border:1px solid {BORDER};border-left:3px solid {sc};margin-bottom:6px;border-radius:4px;">
                    <div style="font-size:0.82rem;font-weight:600;color:{NAVY};margin-bottom:5px;">{r['submarket_name']}</div>
                    <div style="display:grid;grid-template-columns:1fr 1fr;gap:3px;">
                        <div style="font-family:'DM Mono',monospace;font-size:0.62rem;color:{MUTED};">VACANCY</div>
                        <div style="font-family:'DM Mono',monospace;font-size:0.62rem;color:{TEXT};">{r['vacancy']*100:.1f}%</div>
                        <div style="font-family:'DM Mono',monospace;font-size:0.62rem;color:{MUTED};">RENT GROWTH</div>
                        <div style="font-family:'DM Mono',monospace;font-size:0.62rem;color:{RED if r['rent_growth']<0 else GREEN};">{r['rent_growth']*100:+.1f}%</div>
                        <div style="font-family:'DM Mono',monospace;font-size:0.62rem;color:{MUTED};">UNDER CONSTR</div>
                        <div style="font-family:'DM Mono',monospace;font-size:0.62rem;color:{TEXT};">{r['under_constr']:,}</div>
                        <div style="font-family:'DM Mono',monospace;font-size:0.62rem;color:{MUTED};">ABSORPTION</div>
                        <div style="font-family:'DM Mono',monospace;font-size:0.62rem;color:{TEXT};">{r.get('absorption_12mo',0):,}</div>
                        <div style="font-family:'DM Mono',monospace;font-size:0.62rem;color:{MUTED};">CONCESSIONS</div>
                        <div style="font-family:'DM Mono',monospace;font-size:0.62rem;color:{TEXT};">{r.get('concession_pct',0)*100:.1f}%</div>
                        <div style="font-family:'DM Mono',monospace;font-size:0.62rem;color:{MUTED};">SCORE</div>
                        <div style="font-family:'DM Mono',monospace;font-size:0.62rem;color:{sc};font-weight:600;">{r['score']:.0f}/100</div>
                    </div>
                </div>
                """, unsafe_allow_html=True)

    st.markdown("<br>", unsafe_allow_html=True)
    st.markdown('<div class="section-title">Pressure Score Ranked</div>', unsafe_allow_html=True)
    sd = dc.sort_values("score", ascending=True)
    fig7 = go.Figure(go.Bar(
        x=sd["score"], y=sd["submarket_name"], orientation="h",
        marker_color=[sig_color(s) for s in sd["score"]], opacity=0.85,
        text=sd["score"].apply(lambda x: f"{x:.0f}"), textposition="outside",
        textfont=dict(size=9,color=MUTED,family="DM Mono"),
    ))
    fig7.add_vline(x=60, line=dict(color=RED,width=1,dash="dot"), annotation_text="SELL", annotation_font_color=RED)
    fig7.add_vline(x=35, line=dict(color=AMBER,width=1,dash="dot"), annotation_text="HOLD", annotation_font_color=AMBER)
    fig7.update_layout(**PLOTLY_LAYOUT, height=520, xaxis_range=[0,110], yaxis_tickfont_size=10, yaxis_tickfont_family="DM Mono")
    st.plotly_chart(fig7, use_container_width=True)

# ══════════════ TAB 5 ══════════════
with t5:
    fa, fb, fc = st.columns([2, 2, 1])
    with fa:
        search = st.text_input("", placeholder="Search address, project, or ZIP...", label_visibility="collapsed")
    with fb:
        subs = ["All Submarkets"] + sorted(df["submarket_name"].dropna().unique().tolist()) if not df.empty else ["All Submarkets"]
        sub_sel = st.selectbox("", subs, label_visibility="collapsed")
    with fc:
        min_u = st.number_input("", value=5, min_value=5, step=10, label_visibility="collapsed")

    disp = df_f.copy() if not df_f.empty else pd.DataFrame()
    if not disp.empty:
        if search:
            m = (disp["address"].str.contains(search,case=False,na=False) |
                 disp["project_name"].str.contains(search,case=False,na=False) |
                 disp["zip_code"].astype(str).str.contains(search,case=False,na=False))
            disp = disp[m]
        if sub_sel != "All Submarkets":
            disp = disp[disp["submarket_name"] == sub_sel]
        if min_u > 5:
            disp = disp[disp["total_units"] >= min_u]

        st.markdown(f'<div style="font-family:\'DM Mono\',monospace;font-size:0.62rem;color:{MUTED};margin-bottom:0.5rem;">{len(disp):,} permits</div>', unsafe_allow_html=True)
        show = disp[["issue_date","address","zip_code","submarket_name","total_units","project_name","permit_num"]].rename(columns={
            "issue_date":"CO Date","address":"Address","zip_code":"ZIP",
            "submarket_name":"Submarket","total_units":"Units","project_name":"Project","permit_num":"Permit #"
        })
        st.dataframe(show.head(500), use_container_width=True, height=500, hide_index=True)
        st.download_button("Export CSV", disp.to_csv(index=False), f"austin_co_{datetime.now().strftime('%Y%m%d')}.csv", "text/csv")

# ══════════════ TAB 6 — MAP ══════════════
with t6:
    st.markdown('<div class="section-title">Submarket Boundaries & Permit Locations</div>', unsafe_allow_html=True)
    if not df_f.empty:
        map_df = df_f.dropna(subset=["latitude", "longitude"]).copy()
        map_df = map_df[(map_df["latitude"] != 0) & (map_df["longitude"] != 0)]
        if not map_df.empty:
            ma, mb = st.columns([3, 1])
            with mb:
                map_subs = ["All Submarkets"] + sorted(map_df["submarket_name"].dropna().unique().tolist())
                map_sub_sel = st.selectbox("Submarket", map_subs, label_visibility="collapsed", key="map_sub")
                map_min_units = st.slider("Minimum units", 5, 200, 5, key="map_units")
                show_boundaries = st.checkbox("Show submarket boundaries", value=True, key="map_bounds")
            map_show = map_df.copy()
            if map_sub_sel != "All Submarkets":
                map_show = map_show[map_show["submarket_name"] == map_sub_sel]
            if map_min_units > 5:
                map_show = map_show[map_show["total_units"] >= map_min_units]
            with mb:
                st.markdown(f'<div style="font-family:\'DM Mono\',monospace;font-size:0.62rem;color:{MUTED};margin-top:0.5rem;">{len(map_show):,} permits mapped</div>', unsafe_allow_html=True)
            with ma:
                fig_map = go.Figure()

                # Submarket boundary polygons
                if show_boundaries:
                    try:
                        geojson = load_submarket_boundaries()
                        # Color palette for submarkets
                        sm_colors = [
                            "#1a1a2e", "#c8102e", "#2a9d8f", "#e9c46a", "#264653",
                            "#e76f51", "#606c38", "#6d6875", "#0077b6", "#bc6c25",
                            "#457b9d", "#8338ec", "#06d6a0", "#ef476f", "#ffd166",
                            "#118ab2", "#073b4c", "#70a288", "#d4a373", "#588157",
                            "#a7c957", "#6b705c", "#cb997e", "#b5838d", "#e5989b",
                            "#780000", "#023e8a", "#d00000",
                        ]
                        for i, feat in enumerate(geojson["features"]):
                            name = feat["properties"]["submarket_name"]
                            if map_sub_sel != "All Submarkets" and name != map_sub_sel:
                                continue
                            geom = feat["geometry"]
                            color = sm_colors[i % len(sm_colors)]
                            polys = geom.get("coordinates", [])
                            if geom["type"] == "Polygon":
                                polys = [polys]
                            for poly_coords in polys:
                                ring = poly_coords[0]  # exterior ring
                                lons = [c[0] for c in ring]
                                lats = [c[1] for c in ring]
                                fig_map.add_trace(go.Scattermapbox(
                                    lon=lons, lat=lats,
                                    mode="lines",
                                    line=dict(width=2, color=color),
                                    fill="toself",
                                    fillcolor=color.replace(")", ",0.1)").replace("rgb", "rgba") if "rgb" in color else color + "1A",
                                    name=name,
                                    showlegend=(poly_coords == polys[0]),
                                    hoverinfo="name",
                                ))
                    except Exception:
                        pass  # boundaries optional

                # Permit dots
                fig_map.add_trace(go.Scattermapbox(
                    lat=map_show["latitude"], lon=map_show["longitude"],
                    mode="markers",
                    marker=dict(
                        size=map_show["total_units"].clip(5, 300).apply(lambda x: max(4, min(x / 15, 18))),
                        color=ACCENT, opacity=0.7,
                    ),
                    text=map_show.apply(lambda r: f"{r.get('project_name','') or r['address']}<br>{r['total_units']:,} units<br>{r['submarket_name']}", axis=1),
                    hoverinfo="text",
                    name="Permits",
                    showlegend=True,
                ))

                center_lat = map_show["latitude"].mean()
                center_lon = map_show["longitude"].mean()
                fig_map.update_layout(
                    mapbox=dict(
                        style="carto-positron",
                        center=dict(lat=center_lat, lon=center_lon),
                        zoom=9 if map_sub_sel == "All Submarkets" else 11,
                    ),
                    margin=dict(l=0, r=0, t=0, b=0),
                    height=620,
                    legend=dict(
                        bgcolor="rgba(255,255,255,0.85)",
                        font=dict(size=9, family="DM Mono"),
                        x=0.01, y=0.99, xanchor="left", yanchor="top",
                    ),
                    showlegend=True,
                )
                st.plotly_chart(fig_map, use_container_width=True)
        else:
            st.info("No geocoded permits available.")
    else:
        st.info("No permit data loaded.")

# ─────────────────────────────────────────────────────────────────────────────
# POWERPOINT EXPORT — Submarket Report
# ─────────────────────────────────────────────────────────────────────────────
PPTX_NAVY = RGBColor(0x1B, 0x2A, 0x4A)
PPTX_TEAL = RGBColor(0x2A, 0x9D, 0x8F)
PPTX_LTBLUE = RGBColor(0x45, 0x7B, 0x9D)
PPTX_RED = RGBColor(0xC8, 0x10, 0x2E)
PPTX_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PPTX_GRAY = RGBColor(0x6B, 0x72, 0x80)
PPTX_LTGRAY = RGBColor(0xE5, 0xE7, 0xEB)
PPTX_BLACK = RGBColor(0x1F, 0x1F, 0x1F)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
HEADER_H = Inches(0.9)

def _add_text(slide, left, top, width, height, text, font_size=12, color=PPTX_NAVY, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def _add_header_bar(slide, title):
    """Add a dark navy header bar across the top of the slide."""
    from pptx.util import Emu as _Emu
    rect = slide.shapes.add_shape(
        1, Inches(0), Inches(0), SLIDE_W, HEADER_H  # MSO_SHAPE.RECTANGLE = 1
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = PPTX_NAVY
    rect.line.fill.background()
    _add_text(slide, 0.6, 0.15, 10, 0.6, title, 24, PPTX_WHITE, True)

def _add_footer(slide, submarket_name):
    _add_text(slide, 0.6, 6.9, 8, 0.4,
              f"Matthews Real Estate Investment Services  |  {submarket_name}  |  Confidential",
              9, PPTX_GRAY, False, PP_ALIGN.LEFT)
    _add_text(slide, 9, 6.9, 4, 0.4,
              datetime.now().strftime("%B %d, %Y"),
              9, PPTX_GRAY, False, PP_ALIGN.RIGHT)

def _add_table(slide, left, top, width, height, headers, rows, col_widths=None):
    """Add a formatted table to a slide. rows is list of lists of strings."""
    from pptx.table import Table
    n_rows = len(rows) + 1  # +1 for header
    n_cols = len(headers)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top), Inches(width), Inches(height))
    tbl = tbl_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)

    # Header row
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(9)
            paragraph.font.bold = True
            paragraph.font.color.rgb = PPTX_WHITE
            paragraph.font.name = "Calibri"
        cell.fill.solid()
        cell.fill.fore_color.rgb = PPTX_NAVY

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(8)
                paragraph.font.color.rgb = PPTX_BLACK
                paragraph.font.name = "Calibri"
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)

    return tbl_shape

def _render_map_image(map_permits, submarket_name):
    """Render a static map of permits as a PNG bytes buffer using matplotlib + contextily."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    fig, ax = plt.subplots(1, 1, figsize=(10, 7))

    # Draw submarket boundary if available
    try:
        conn = psycopg2.connect(DB_DSN)
        cur = conn.cursor()
        cur.execute("""
            SELECT ST_AsGeoJSON(geom)::text FROM costar_submarkets
            WHERE submarket_name = %s AND geom IS NOT NULL
        """, (submarket_name,))
        row = cur.fetchone()
        conn.close()
        if row:
            geom = json.loads(row[0])
            from shapely.geometry import shape as shp_shape
            boundary = shp_shape(geom)
            if boundary.geom_type == "MultiPolygon":
                for poly in boundary.geoms:
                    xs, ys = poly.exterior.xy
                    ax.plot(xs, ys, color="#1B2A4A", linewidth=2, zorder=2)
                    ax.fill(xs, ys, alpha=0.08, color="#1B2A4A", zorder=1)
            else:
                xs, ys = boundary.exterior.xy
                ax.plot(xs, ys, color="#1B2A4A", linewidth=2, zorder=2)
                ax.fill(xs, ys, alpha=0.08, color="#1B2A4A", zorder=1)
    except Exception:
        pass

    if not map_permits.empty:
        geo = map_permits.dropna(subset=["latitude", "longitude"])
        geo = geo[(geo["latitude"] != 0) & (geo["longitude"] != 0)]
        if not geo.empty:
            sizes = geo["total_units"].clip(10, 500).values * 0.5
            colors = geo["delivery_year"].values
            sc = ax.scatter(geo["longitude"], geo["latitude"],
                           c=colors, cmap="YlGnBu", s=sizes,
                           alpha=0.75, edgecolors="#1B2A4A", linewidth=0.5, zorder=3)
            cbar = plt.colorbar(sc, ax=ax, shrink=0.6, pad=0.02)
            cbar.set_label("Delivery Year", fontsize=8)
            cbar.ax.tick_params(labelsize=7)

            # Add basemap tiles
            try:
                import contextily as ctx
                pad = 0.01
                ax.set_xlim(geo["longitude"].min() - pad, geo["longitude"].max() + pad)
                ax.set_ylim(geo["latitude"].min() - pad, geo["latitude"].max() + pad)
                ctx.add_basemap(ax, crs="EPSG:4326", source=ctx.providers.CartoDB.Positron, zoom=12)
            except Exception:
                pass

    ax.set_title(f"{submarket_name} — Permit Locations", fontsize=12, fontweight="bold", color="#1B2A4A")
    ax.set_xlabel("")
    ax.set_ylabel("")
    ax.tick_params(labelsize=7)
    fig.tight_layout()

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf

def _chart_to_image(fig_plotly, width=900, height=500):
    """Convert a plotly figure to PNG bytes using kaleido/orca, fallback to matplotlib."""
    try:
        buf = io.BytesIO()
        fig_plotly.write_image(buf, format="png", width=width, height=height, scale=2)
        buf.seek(0)
        return buf
    except Exception:
        return None

def build_submarket_pptx(submarket_name, dc_df, df_all, dq_all):
    """Build a full submarket report PowerPoint deck."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Filter data to the selected submarket
    sm_permits = df_all[df_all["submarket_name"] == submarket_name].copy() if not df_all.empty else pd.DataFrame()
    sm_costar = dc_df[dc_df["submarket_name"] == submarket_name].iloc[0] if submarket_name in dc_df["submarket_name"].values else None

    # Metro-wide delivery by year (for comparison)
    metro_by_year = df_all.groupby("delivery_year")["total_units"].sum().reset_index() if not df_all.empty else pd.DataFrame()

    # Submarket delivery by year
    sm_by_year = sm_permits.groupby("delivery_year")["total_units"].sum().reset_index() if not sm_permits.empty else pd.DataFrame()

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE 1: TITLE
    # ═══════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PPTX_NAVY
    _add_text(slide, 0.8, 1.2, 11, 0.8, "AUSTIN MULTIFAMILY INTELLIGENCE", 20, PPTX_LTGRAY, False)
    _add_text(slide, 0.8, 2.0, 11, 1.5, submarket_name.upper(), 48, PPTX_WHITE, True)
    _add_text(slide, 0.8, 3.8, 8, 0.8, "Submarket Report", 26, PPTX_TEAL, False)
    _add_text(slide, 0.8, 5.2, 8, 0.6,
              f"Matthews Real Estate Investment Services  |  {datetime.now().strftime('%B %d, %Y')}",
              14, PPTX_GRAY, False)

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE 2: SUPPLY PIPELINE SUMMARY
    # ═══════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, f"SUPPLY PIPELINE — {submarket_name.upper()}")

    # KPI boxes
    kpi_data = []
    if sm_costar is not None:
        kpi_data = [
            ("Total Inventory", f"{int(sm_costar['inventory']):,}"),
            ("Under Construction", f"{int(sm_costar['under_constr']):,}"),
            ("Delivered (12mo)", f"{int(sm_costar['delivered_12mo']):,}"),
            ("Net Absorption (12mo)", f"{int(sm_costar.get('absorption_12mo', 0)):,}"),
            ("Vacancy Rate", f"{sm_costar['vacancy']*100:.1f}%"),
            ("Rent Growth", f"{sm_costar['rent_growth']*100:+.1f}%"),
        ]
    else:
        total_units = int(sm_permits["total_units"].sum()) if not sm_permits.empty else 0
        kpi_data = [
            ("Total CO Units", f"{total_units:,}"),
            ("Projects", f"{len(sm_permits):,}"),
            ("Avg Project Size", f"{int(sm_permits['total_units'].mean()):,}" if not sm_permits.empty else "N/A"),
        ]

    for i, (label, val) in enumerate(kpi_data):
        col = i % 3
        row_idx = i // 3
        x = 0.6 + col * 4.1
        y_base = 1.2 + row_idx * 1.4
        # KPI card background
        rect = slide.shapes.add_shape(1, Inches(x), Inches(y_base), Inches(3.8), Inches(1.1))
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
        rect.line.color.rgb = PPTX_LTGRAY
        rect.line.width = Pt(1)
        _add_text(slide, x + 0.15, y_base + 0.1, 3.5, 0.3, label, 10, PPTX_GRAY, False)
        _add_text(slide, x + 0.15, y_base + 0.4, 3.5, 0.6, val, 28, PPTX_NAVY, True)

    # Delivery by year bar chart (inline as table-bars since kaleido may not be available)
    if not sm_by_year.empty:
        _add_text(slide, 0.6, 4.0, 6, 0.4, "UNITS DELIVERED BY YEAR", 12, PPTX_NAVY, True)
        chart_years = sm_by_year.sort_values("delivery_year").tail(15)
        chart_rows = []
        max_units = chart_years["total_units"].max() if not chart_years.empty else 1
        for _, r in chart_years.iterrows():
            bar_len = int(r["total_units"] / max_units * 30) if max_units > 0 else 0
            bar = "█" * bar_len
            chart_rows.append([str(int(r["delivery_year"])), f"{int(r['total_units']):,}", bar])
        _add_table(slide, 0.6, 4.4, 7.5, min(3.0, len(chart_rows) * 0.2 + 0.3),
                   ["Year", "Units", ""], chart_rows,
                   col_widths=[0.8, 1.0, 5.7])

    # Top 10 projects table
    if not sm_permits.empty:
        top10 = sm_permits.nlargest(10, "total_units")
        _add_text(slide, 8.4, 1.2, 4.5, 0.4, "TOP 10 PROJECTS", 12, PPTX_NAVY, True)
        t10_rows = []
        for _, r in top10.iterrows():
            addr = str(r.get("address", ""))[:35]
            t10_rows.append([
                addr,
                f"{int(r['total_units']):,}",
                str(r.get("issue_date", ""))[:10],
            ])
        _add_table(slide, 8.4, 1.6, 4.5, min(3.5, len(t10_rows) * 0.28 + 0.3),
                   ["Address", "Units", "CO Date"], t10_rows,
                   col_widths=[2.5, 0.8, 1.2])

    _add_footer(slide, submarket_name)

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE 3+: PERMIT BROWSER (paginated)
    # ═══════════════════════════════════════════════════════════════════════
    if not sm_permits.empty:
        sorted_permits = sm_permits.sort_values("issue_date", ascending=False)
        total_permits = len(sorted_permits)
        total_sm_units = int(sorted_permits["total_units"].sum())
        avg_size = int(sorted_permits["total_units"].mean())
        rows_per_page = 25
        pages = (total_permits + rows_per_page - 1) // rows_per_page

        for page in range(pages):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            page_label = f" (Page {page+1}/{pages})" if pages > 1 else ""
            _add_header_bar(slide, f"PERMIT BROWSER — {submarket_name.upper()}{page_label}")

            if page == 0:
                # Summary stats on first page
                stats_text = f"{total_permits:,} permits  |  {total_sm_units:,} total units  |  {avg_size:,} avg project size"
                _add_text(slide, 0.6, 1.1, 10, 0.3, stats_text, 11, PPTX_TEAL, True)
                tbl_top = 1.5
            else:
                tbl_top = 1.2

            chunk = sorted_permits.iloc[page * rows_per_page : (page + 1) * rows_per_page]
            permit_rows = []
            for _, r in chunk.iterrows():
                permit_rows.append([
                    str(r.get("issue_date", ""))[:10],
                    str(r.get("address", ""))[:40],
                    str(r.get("zip_code", "")),
                    f"{int(r['total_units']):,}",
                    str(r.get("project_name", "") or "")[:30],
                    str(r.get("permit_num", "")),
                ])

            _add_table(slide, 0.6, tbl_top, 12.1, min(5.5, len(permit_rows) * 0.2 + 0.3),
                       ["CO Date", "Address", "ZIP", "Units", "Project Description", "Permit #"],
                       permit_rows,
                       col_widths=[1.2, 3.5, 0.8, 0.8, 3.0, 2.8])
            _add_footer(slide, submarket_name)

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE: MAP
    # ═══════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, f"PERMIT MAP — {submarket_name.upper()}")

    try:
        map_buf = _render_map_image(sm_permits, submarket_name)
        slide.shapes.add_picture(map_buf, Inches(0.6), Inches(1.1), Inches(12.1), Inches(5.6))
    except Exception as e:
        _add_text(slide, 2, 3, 8, 1, f"Map rendering unavailable: {e}", 14, PPTX_GRAY)

    _add_footer(slide, submarket_name)

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE: DELIVERY TREND (submarket vs metro)
    # ═══════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, f"DELIVERY TREND — {submarket_name.upper()}")

    if not sm_by_year.empty:
        # Build a text-based chart since we can't guarantee plotly image export
        merged = sm_by_year.rename(columns={"total_units": "sm_units"}).merge(
            metro_by_year.rename(columns={"total_units": "metro_units"}),
            on="delivery_year", how="outer"
        ).fillna(0).sort_values("delivery_year")
        merged = merged[merged["delivery_year"] >= 1990]  # reasonable range

        _add_text(slide, 0.6, 1.1, 5, 0.4, "Units delivered per year — submarket vs metro average", 11, PPTX_GRAY)

        # Submarket trend as table with visual bars
        chart_rows = []
        max_sm = merged["sm_units"].max() if not merged.empty else 1
        n_submarkets = dc_df["submarket_name"].nunique()
        for _, r in merged.iterrows():
            yr = int(r["delivery_year"])
            sm_u = int(r["sm_units"])
            metro_avg = int(r["metro_units"] / n_submarkets) if n_submarkets > 0 else 0
            bar_sm = "█" * int(sm_u / max(max_sm, 1) * 25)
            chart_rows.append([str(yr), f"{sm_u:,}", f"{metro_avg:,}", bar_sm])

        _add_table(slide, 0.6, 1.5, 12.1, min(5.2, len(chart_rows) * 0.18 + 0.3),
                   ["Year", f"{submarket_name} Units", "Metro Avg/Submarket", "Submarket Volume"],
                   chart_rows,
                   col_widths=[1.0, 2.0, 2.5, 6.6])
    else:
        _add_text(slide, 2, 3, 8, 1, "No delivery data available for this submarket.", 14, PPTX_GRAY)

    _add_footer(slide, submarket_name)

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE: ACTIVE CONSTRUCTION (recent 2 years as proxy)
    # ═══════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, f"RECENT ACTIVITY — {submarket_name.upper()}")

    current_year = datetime.now().year
    if not sm_permits.empty:
        recent = sm_permits[sm_permits["delivery_year"] >= current_year - 2].sort_values("issue_date", ascending=False)
        _add_text(slide, 0.6, 1.1, 10, 0.3,
                  f"Projects with CO issued {current_year - 2}–{current_year}  |  {len(recent):,} projects  |  {int(recent['total_units'].sum()):,} units",
                  11, PPTX_TEAL, True)

        if not recent.empty:
            recent_rows = []
            for _, r in recent.head(30).iterrows():
                recent_rows.append([
                    str(r.get("issue_date", ""))[:10],
                    str(r.get("address", ""))[:40],
                    f"{int(r['total_units']):,}",
                    str(r.get("project_name", "") or "")[:35],
                    str(r.get("permit_num", "")),
                ])
            _add_table(slide, 0.6, 1.5, 12.1, min(5.2, len(recent_rows) * 0.2 + 0.3),
                       ["CO Date", "Address", "Units", "Project", "Permit #"],
                       recent_rows,
                       col_widths=[1.2, 3.8, 1.0, 3.3, 2.8])
        else:
            _add_text(slide, 2, 3, 8, 1, "No recent projects in this submarket.", 14, PPTX_GRAY)
    else:
        _add_text(slide, 2, 3, 8, 1, "No permit data available.", 14, PPTX_GRAY)

    _add_footer(slide, submarket_name)

    # ═══════════════════════════════════════════════════════════════════════
    # SLIDE: METHODOLOGY
    # ═══════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_bar(slide, "METHODOLOGY & DATA SOURCES")
    methodology = (
        "Supply Pressure Score (0-100)\n"
        "Composite of seven weighted factors:\n"
        "  • Vacancy Rate: 25 pts\n"
        "  • 12-Month Deliveries vs Inventory: 20 pts\n"
        "  • Under Construction vs Inventory: 20 pts\n"
        "  • Rent Growth (inverted): 15 pts\n"
        "  • Absorption vs Deliveries: 10 pts\n"
        "  • Avg Days on Market: 5 pts\n"
        "  • Concession Rate: 5 pts\n\n"
        "Investment Signals: BUY (<35) | HOLD (35-59) | SELL (60+)\n\n"
        "Data Sources:\n"
        "  • City of Austin Open Data Portal — C-104, C-105, C-106 Certificates of Occupancy\n"
        "  • CoStar Group — vacancy, rent, absorption, pipeline metrics\n"
        "  • Census TIGER/Line ZCTA — submarket boundary polygons\n"
        "  • PostGIS ST_Contains — point-in-polygon submarket assignment (98.6% match rate)\n"
        "  • Filtered to NEW permits, 5–1,000 units, deduplicated by master permit number"
    )
    _add_text(slide, 0.6, 1.3, 11.5, 5.0, methodology, 13, PPTX_BLACK)
    _add_footer(slide, submarket_name)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf

with st.sidebar:
    st.markdown(f'<div style="font-family:\'DM Mono\',monospace;font-size:0.65rem;color:{MUTED};letter-spacing:0.15em;text-transform:uppercase;margin-top:1.5rem;margin-bottom:0.5rem;">Export</div>', unsafe_allow_html=True)

    # Submarket selector for export
    export_subs = sorted(df["submarket_name"].dropna().unique().tolist()) if not df.empty else []
    export_sub = st.selectbox("Submarket report", ["All Submarkets"] + export_subs,
                              label_visibility="collapsed", key="export_sub")

    if export_sub != "All Submarkets":
        pptx_buf = build_submarket_pptx(export_sub, dc, df, dq)
        fname = export_sub.lower().replace(" ", "_")
        st.download_button(
            label=f"Export {export_sub} Report",
            data=pptx_buf,
            file_name=f"{fname}_report_{datetime.now().strftime('%Y%m%d')}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    else:
        # Metro-wide summary export (legacy)
        pptx_buf = build_submarket_pptx("All Submarkets", dc, df, dq)
        st.download_button(
            label="Export Metro Summary",
            data=pptx_buf,
            file_name=f"austin_mf_intelligence_{datetime.now().strftime('%Y%m%d')}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

# ─────────────────────────────────────────────────────────────────────────────
# FOOTER
# ─────────────────────────────────────────────────────────────────────────────
st.markdown(f"""
<div style="margin-top:3rem;padding-top:1rem;border-top:1px solid {BORDER};display:flex;justify-content:space-between;align-items:center;">
    <div style="font-family:'DM Mono',monospace;font-size:0.6rem;color:{MUTED};">DATA: Austin Open Data Portal · CoStar Group · Certificates of Occupancy</div>
    <div style="font-family:'DM Mono',monospace;font-size:0.6rem;color:{MUTED};letter-spacing:0.12em;">MATTHEWS REAL ESTATE INVESTMENT SERVICES</div>
</div>
""", unsafe_allow_html=True)
